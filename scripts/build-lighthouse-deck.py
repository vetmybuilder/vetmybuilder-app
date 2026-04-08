#!/usr/bin/env python3
"""
build-lighthouse-deck.py

Builds the Project Lighthouse internal planning deck — the AI Phase 2
roadmap for VetMyBuilder.

This is a *reference document* for the founder, not an investor pitch.
It captures the full AI strategy: what to ship, when, what data to
collect, what synthetic data is good for, and the week-1 build plan.

Companion to scripts/build-pitch-deck.py — same visual style so the
two decks read as a series.

Usage:
    python3 scripts/build-lighthouse-deck.py
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── Paths ────────────────────────────────────────────────────────────

ROOT = Path(__file__).resolve().parents[1]
DECK_PATH = ROOT / "project-lighthouse" / "project-lighthouse-plan.pptx"

# ── Theme (matches the investor pitch deck) ─────────────────────────

BG = RGBColor(0x0F, 0x17, 0x2A)        # slate-900
PANEL = RGBColor(0x1E, 0x29, 0x3B)     # slate-800
PANEL_LIGHT = RGBColor(0x33, 0x40, 0x55)  # slate-700
TEXT = RGBColor(0xF8, 0xFA, 0xFC)      # slate-50
DIM = RGBColor(0x94, 0xA3, 0xB8)       # slate-400
DIMMER = RGBColor(0x64, 0x74, 0x8B)    # slate-500
ACCENT = RGBColor(0xEF, 0x44, 0x44)    # red-500
ACCENT_DIM = RGBColor(0xFB, 0xA1, 0xA1)  # red-300
GREEN = RGBColor(0x10, 0xB9, 0x81)     # emerald-500
AMBER = RGBColor(0xF5, 0x9E, 0x0B)     # amber-500

FONT = "Inter"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
TOTAL_SLIDES = 18

# ── Slide-builder helpers ────────────────────────────────────────────


def find_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in (layout.name or "").lower():
            return layout
    return prs.slide_layouts[6]


def add_slide(prs: Presentation):
    slide = prs.slides.add_slide(find_blank_layout(prs))
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.15,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_bullets(
    slide,
    *,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 14,
    color: RGBColor = TEXT,
    line_spacing: float = 1.3,
    bullet_char: str = "•",
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_eyebrow(slide, text: str, *, top: float = 0.55, color: RGBColor = ACCENT):
    add_text(
        slide,
        text=text.upper(),
        left=0.7,
        top=top,
        width=12,
        height=0.4,
        size=12,
        bold=True,
        color=color,
    )


def add_footer(slide, page: int, total: int = TOTAL_SLIDES):
    add_text(
        slide,
        text="Project Lighthouse  ·  AI Phase 2",
        left=0.7,
        top=SLIDE_H_IN - 0.55,
        width=8,
        height=0.3,
        size=10,
        color=DIM,
    )
    add_text(
        slide,
        text=f"{page} / {total}",
        left=SLIDE_W_IN - 1.3,
        top=SLIDE_H_IN - 0.55,
        width=1,
        height=0.3,
        size=10,
        color=DIM,
        align=PP_ALIGN.RIGHT,
    )


def panel(slide, *, left: float, top: float, width: float, height: float,
          fill: RGBColor = PANEL, border: RGBColor = ACCENT,
          border_width: float = 1.0, rounded: bool = True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is not None:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ─────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────


def slide_01_title(prs):
    slide = add_slide(prs)

    add_text(
        slide,
        text="Project Lighthouse",
        left=0.7, top=2.0, width=12, height=1.4,
        size=72, bold=True, color=TEXT,
    )
    add_text(
        slide,
        text="The AI Phase 2 plan for VetMyBuilder",
        left=0.7, top=3.35, width=12, height=0.7,
        size=24, color=ACCENT_DIM,
    )
    add_text(
        slide,
        text=(
            "An internal planning document covering what AI to build, when, what data to "
            "capture now, and how to graduate from synthetic to real training data."
        ),
        left=0.7, top=4.25, width=12, height=1.0,
        size=15, color=DIM,
    )

    # Status pill
    pill = panel(
        slide, left=0.7, top=5.35, width=2.6, height=0.5,
        fill=PANEL, border=ACCENT, border_width=1.5,
    )
    add_text(
        slide,
        text="STATUS · DRAFT",
        left=0.7, top=5.43, width=2.6, height=0.35,
        size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    add_text(
        slide,
        text="Companion to vetmybuilder-pitch-deck-2026.pptx · Internal use only",
        left=0.7, top=SLIDE_H_IN - 0.65, width=12, height=0.4,
        size=11, color=DIM,
    )


def slide_02_tldr(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "TL;DR")
    add_text(
        slide,
        text="AI is how VetMyBuilder becomes\na service, not a list.",
        left=0.7, top=1.0, width=12, height=2.0,
        size=40, bold=True,
    )

    points = [
        ("The bet",
         "Most marketplaces are search engines for suppliers. AI lets us be a "
         "decision assistant for the homeowner — the difference between Yelp "
         "and Booking.com."),
        ("The wedge",
         "Our recommendation graph (neighbours endorsing builders they actually "
         "used) is unique training data. Nobody else can replicate it without "
         "running the same marketplace."),
        ("The plan",
         "Ship 4 LLM-driven AI features in 4 weeks (Tier 1). Collect the right "
         "data from day one. Train real models in months 3–6 (Tier 2). The moat "
         "compounds quietly while we focus on growth."),
    ]
    y = 3.4
    for head, body in points:
        add_text(slide, text=head,
                 left=0.7, top=y, width=12, height=0.4,
                 size=15, bold=True, color=ACCENT_DIM)
        add_text(slide, text=body,
                 left=0.7, top=y + 0.4, width=12, height=0.85,
                 size=13, color=DIM, line_spacing=1.35)
        y += 1.25

    add_footer(slide, 2)


def slide_03_why_now(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Why now")
    add_text(
        slide,
        text="The four conditions are aligned.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=32, bold=True,
    )

    cards = [
        ("01", "Phase 1 is shipped",
         "The core marketplace is live. Hires, recommendations, scoring, trust "
         "signals — all in production. The substrate exists."),
        ("02", "Recommendation graph is unique",
         "Nobody else has a marketplace built on neighbour-endorsed builders. "
         "Every project we close is training data no competitor can buy."),
        ("03", "AI APIs got cheap",
         "Claude Haiku and gpt-4o-mini cost cents per thousand calls. A solo "
         "founder can ship production AI features without an ML team."),
        ("04", "Slot in the pitch story",
         "Slide 8 of the pitch deck (Trust & defensibility) becomes much "
         "stronger with a credible AI moat already in production."),
    ]

    card_w = 2.95
    card_h = 3.7
    gap = 0.18
    start_left = (SLIDE_W_IN - (card_w * 4 + gap * 3)) / 2
    top = 2.5

    for i, (num, head, body) in enumerate(cards):
        left = start_left + i * (card_w + gap)
        panel(slide, left=left, top=top, width=card_w, height=card_h)
        add_text(slide, text=num,
                 left=left + 0.3, top=top + 0.3, width=card_w - 0.6, height=0.5,
                 size=22, bold=True, color=ACCENT)
        add_text(slide, text=head,
                 left=left + 0.3, top=top + 0.85, width=card_w - 0.6, height=0.7,
                 size=15, bold=True, color=TEXT)
        add_text(slide, text=body,
                 left=left + 0.3, top=top + 1.55, width=card_w - 0.6, height=2.0,
                 size=11, color=DIM, line_spacing=1.4)

    add_footer(slide, 3)


def slide_04_three_flywheels(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "The moat")
    add_text(
        slide,
        text="Three compounding flywheels — only AI unlocks them.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    flywheels = [
        ("Project understanding compounds",
         "More projects → more (project, hired-builder) pairs → better re-ranker → "
         "better matches → more hires → more projects.",
         "After 1k projects: train the classifier.\nAfter 10k: personalise per project type.\nAfter 100k: per project × postcode × property."),
        ("Trust scores get more accurate",
         "More moderated abuse → labelled fraud examples → better classifier → "
         "catches novel attacks before competitors do.",
         "Bootstrap with simulated abuse.\nReplace with real labels at month 3.\nAuto-improves indefinitely."),
        ("Match outcomes train the matcher",
         "Every closed hire is a label. The match engine learns which builder profiles "
         "win which kinds of projects in which postcodes.",
         "Needs ~500 closed hires to train.\nGrows monotonically with usage.\nUnfakeable by competitors."),
    ]

    y = 2.4
    for head, body, when in flywheels:
        panel(slide, left=0.7, top=y, width=12.0, height=1.4)
        add_text(slide, text=head,
                 left=0.95, top=y + 0.12, width=8.0, height=0.35,
                 size=14, bold=True, color=ACCENT_DIM)
        add_text(slide, text=body,
                 left=0.95, top=y + 0.5, width=8.0, height=0.85,
                 size=11, color=TEXT, line_spacing=1.4)
        add_text(slide, text=when,
                 left=9.1, top=y + 0.18, width=3.5, height=1.15,
                 size=10, color=DIM, line_spacing=1.4)
        y += 1.55

    add_footer(slide, 4)


def slide_05_capability_map(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Capability map")
    add_text(
        slide,
        text="Twelve AI surfaces, six capability domains.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    domains = [
        ("Project understanding",
         ["Free-text → structured tags",
          "Suggested budget band",
          "Suggested complexity / urgency"]),
        ("Match scoring",
         ["Semantic embedding match",
          "Re-ranker on hire outcomes",
          "Personalised per homeowner"]),
        ("Trust & fraud",
         ["LLM content moderation",
          "Behavioural cluster detection",
          "Anomaly outlier scoring"]),
        ("Review summarisation",
         ["3-bullet builder summaries",
          "Theme extraction",
          "Sentiment + reliability scoring"]),
        ("Builder coach",
         ["Profile gap detection",
          "Weekly growth report",
          "Predictive coaching"]),
        ("Decision assistant",
         ["Side-by-side comparison",
          "Recommendation w/ rationale",
          "Conversational project setup"]),
    ]

    cols = 3
    card_w = 4.0
    card_h = 2.4
    gap_x = 0.2
    gap_y = 0.2
    start_left = (SLIDE_W_IN - (card_w * cols + gap_x * (cols - 1))) / 2
    top = 2.4

    for i, (name, items) in enumerate(domains):
        col = i % cols
        row = i // cols
        left = start_left + col * (card_w + gap_x)
        y = top + row * (card_h + gap_y)
        panel(slide, left=left, top=y, width=card_w, height=card_h)
        add_text(slide, text=name,
                 left=left + 0.3, top=y + 0.2, width=card_w - 0.6, height=0.4,
                 size=14, bold=True, color=ACCENT_DIM)
        add_bullets(slide,
                    bullets=items,
                    left=left + 0.3, top=y + 0.7,
                    width=card_w - 0.6, height=card_h - 0.9,
                    size=11, color=TEXT, line_spacing=1.3,
                    bullet_char="·")

    add_footer(slide, 5)


def slide_06_tier1_mvp(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Tier 1 — MVP", color=GREEN)
    add_text(
        slide,
        text="Ship in the next 4 weeks. No data threshold.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    features = [
        ("Project classifier",
         "LLM call at submission",
         "~£0.005 per project",
         "Every other AI feature depends on this. Single Anthropic call returns "
         "structured JSON: type, scope, complexity, urgency, recommended trades, "
         "price band. Stored in project_classifications table."),
        ("Recommendation summariser",
         "Nightly LLM cron",
         "~£5 per 1k builders / month",
         "For each builder with ≥3 recommendations, generate a 3-bullet summary "
         "of what reviewers consistently say. Cached. Refreshes only when content "
         "changes. Displayed on builder profile."),
        ("Profile coach (rule-based)",
         "SQL-only, no LLM",
         "Free",
         "Gap-detection on the tradesman dashboard. 'Add 3 more photos.' "
         "'Verify with Companies House.' All deterministic, all from existing "
         "schema. Builds the habit of in-app coaching."),
        ("Smart job match badges",
         "Rule-based score",
         "Free",
         "On the tradesman jobs list, score each project against the builder's "
         "profile (trade match, area match, scope match) and surface a 'good fit' "
         "badge. Uses the structured project from feature 1."),
    ]

    y = 2.45
    for head, mech, cost, body in features:
        panel(slide, left=0.7, top=y, width=12.0, height=1.05)
        add_text(slide, text=head,
                 left=0.95, top=y + 0.12, width=5.0, height=0.4,
                 size=14, bold=True, color=ACCENT_DIM)
        add_text(slide, text=f"{mech}  ·  {cost}",
                 left=0.95, top=y + 0.5, width=5.0, height=0.35,
                 size=10, color=DIMMER)
        add_text(slide, text=body,
                 left=6.1, top=y + 0.12, width=6.6, height=0.85,
                 size=10, color=TEXT, line_spacing=1.35)
        y += 1.15

    # Total cost callout at bottom
    panel(slide, left=0.7, top=7.0 - 0.45, width=12, height=0.0)  # placeholder
    add_text(slide,
             text="Total monthly cost at MVP traffic: ~£30.   Engineering time: ~4 weeks of solo dev.",
             left=0.7, top=7.05 - 0.55, width=12, height=0.4,
             size=11, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 6)


def slide_07_tier2_data(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Tier 2 — Once data accumulates", color=AMBER)
    add_text(
        slide,
        text="Months 3–6. Real models, real outcomes.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    features = [
        ("Trust / fraud classifier",
         "XGBoost on labelled abuse signals",
         "Threshold: ~3 months of moderated flags",
         "Replaces hand-coded rules. Trained on (account features → fraud "
         "probability). Catches combinations rules can't anticipate. "
         "Bootstrap with simulated fraud — see slide 13."),
        ("Match engine v2",
         "Logistic regression on hire outcomes",
         "Threshold: ~500 closed hires",
         "Re-ranks the top recommendations card by predicted hire probability. "
         "Stays interpretable. Built on top of the deterministic computeScore "
         "as a learned correction layer."),
        ("Decision Assistant",
         "LLM with structured context",
         "Threshold: working summariser + match engine",
         "When a homeowner has 3+ shortlisted builders, an LLM generates a "
         "side-by-side comparison and recommends one with rationale. The "
         "flagship feature — turns the marketplace into a service."),
    ]

    y = 2.5
    for head, mech, threshold, body in features:
        panel(slide, left=0.7, top=y, width=12.0, height=1.4)
        add_text(slide, text=head,
                 left=0.95, top=y + 0.15, width=5.0, height=0.4,
                 size=14, bold=True, color=ACCENT_DIM)
        add_text(slide, text=mech,
                 left=0.95, top=y + 0.55, width=5.0, height=0.3,
                 size=10, color=DIMMER)
        add_text(slide, text=threshold,
                 left=0.95, top=y + 0.85, width=5.0, height=0.3,
                 size=10, color=AMBER)
        add_text(slide, text=body,
                 left=6.1, top=y + 0.15, width=6.6, height=1.15,
                 size=10, color=TEXT, line_spacing=1.4)
        y += 1.55

    add_footer(slide, 7)


def slide_08_tier3_scale(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Tier 3 — Build at scale", color=DIM)
    add_text(
        slide,
        text="12+ months. The features that need real volume.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    features = [
        ("Personalised match models",
         "A per-homeowner preference vector learned from browsing and hire history. "
         "Different people care about different things; the model learns it."),
        ("Anomaly detection for fraud",
         "Unsupervised model that flags accounts statistically unlike anything else "
         "on the platform — catches abuse before you've seen it."),
        ("Automated project briefs",
         "Free-text description → 1-page brief with materials, expected price range, "
         "common pitfalls, questions to ask each quote, red flags to watch for. "
         "The feature that makes VMB feel like a friend in the trade."),
        ("Tradesman growth coach v2",
         "Predictive coaching: 'builders who added photos AND verified AND responded "
         "within 2h had 70% activation. You've done 2 of 3.' Needs the full hire "
         "outcome dataset."),
    ]

    y = 2.4
    for head, body in features:
        panel(slide, left=0.7, top=y, width=12.0, height=1.0,
              border=DIMMER)
        add_text(slide, text=head,
                 left=0.95, top=y + 0.12, width=11.0, height=0.4,
                 size=14, bold=True, color=DIM)
        add_text(slide, text=body,
                 left=0.95, top=y + 0.5, width=11.0, height=0.5,
                 size=11, color=DIMMER, line_spacing=1.4)
        y += 1.15

    add_footer(slide, 8)


def slide_09_data_strategy(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Data strategy")
    add_text(
        slide,
        text="Capture now. Use later. Don't act on it yet.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )
    add_text(
        slide,
        text=(
            "The single biggest mistake solo founders make with ML: collecting "
            "the wrong data. Add these four tables in week 1 — even if you "
            "don't use them for months."
        ),
        left=0.7, top=1.95, width=12, height=0.6,
        size=12, color=DIM,
    )

    tables = [
        ("project_classifications",
         "Raw description + structured JSON\n+ classifier_version + cost_pence",
         "Powers the matcher and decision assistant"),
        ("match_observations",
         "Every (project, builder) pair shown\n+ rank + action + outcome",
         "THE gold dataset for training the re-ranker"),
        ("recommendation_signals",
         "Word count, sentiment, generic_score\n+ themes JSON",
         "Feeds the summariser and trust score"),
        ("ai_inference_log",
         "Feature, model, prompt, response\n+ cost_pence + latency_ms",
         "Audit trail + retraining + cost monitoring"),
    ]

    card_w = 2.95
    card_h = 3.4
    gap = 0.18
    start_left = (SLIDE_W_IN - (card_w * 4 + gap * 3)) / 2
    top = 2.85

    for i, (name, schema, why) in enumerate(tables):
        left = start_left + i * (card_w + gap)
        panel(slide, left=left, top=top, width=card_w, height=card_h)
        add_text(slide, text=name,
                 left=left + 0.25, top=top + 0.3, width=card_w - 0.5, height=0.4,
                 size=13, bold=True, color=ACCENT)
        add_text(slide, text=schema,
                 left=left + 0.25, top=top + 0.85, width=card_w - 0.5, height=1.4,
                 size=10, color=TEXT, line_spacing=1.4)
        # Divider
        add_text(slide, text="WHY",
                 left=left + 0.25, top=top + 2.3, width=card_w - 0.5, height=0.3,
                 size=8, bold=True, color=DIMMER)
        add_text(slide, text=why,
                 left=left + 0.25, top=top + 2.55, width=card_w - 0.5, height=0.7,
                 size=10, color=DIM, line_spacing=1.35)

    # Most important callout
    add_text(slide,
             text="match_observations is the single most valuable table — it's the dataset every future model is trained on.",
             left=0.7, top=6.45, width=12, height=0.4,
             size=11, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 9)


def slide_10_stack_fit(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Stack integration")
    add_text(
        slide,
        text="No new infrastructure. Drop into the existing stack.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )

    layers = [
        ("Next.js frontend",
         "New React components. Each AI feature surfaces as a card or badge. "
         "No state management changes, no new auth.",
         "/projects/:id    summary card\n/builders/:id    What people say\n/tradesman/projects    fit badges\n/projects/:id/compare    Decision Assistant"),
        ("Express API",
         "New /api/ai/* namespace. Each endpoint is a thin wrapper around "
         "an LLM call. Auth-gated through existing middleware.",
         "POST /api/ai/projects/:id/classify\nGET  /api/ai/builders/:id/summary\nGET  /api/ai/tradesmen/me/coach\nPOST /api/ai/admin/cron/refresh-summaries"),
        ("MySQL",
         "Four new tables (slide 9). No schema changes to existing tables. "
         "Vector storage as JSON columns — fine up to ~10k builders.",
         "project_classifications\nmatch_observations\nrecommendation_signals\nai_inference_log"),
        ("LLM provider",
         "Anthropic Claude Haiku as default (cheap + structured output). "
         "OpenAI gpt-4o-mini as fallback. Local Llama via Ollama for dev.",
         "claude-haiku-4-5\ngpt-4o-mini\nllama-3.1-8b (local)"),
    ]

    y = 2.35
    for layer, body, examples in layers:
        panel(slide, left=0.7, top=y, width=12.0, height=1.05)
        add_text(slide, text=layer,
                 left=0.95, top=y + 0.13, width=3.5, height=0.4,
                 size=13, bold=True, color=ACCENT_DIM)
        add_text(slide, text=body,
                 left=0.95, top=y + 0.5, width=3.5, height=0.5,
                 size=10, color=TEXT, line_spacing=1.35)
        add_text(slide, text=examples,
                 left=4.7, top=y + 0.13, width=8.0, height=0.9,
                 size=10, color=DIM, line_spacing=1.3)
        y += 1.15

    add_footer(slide, 10)


def slide_11_cost_model(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Cost model")
    add_text(
        slide,
        text="AI is a budget rounding error at MVP. Plan for scale.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )

    # 3 cost tiles for different volume bands
    bands = [
        ("MVP", "100 projects / month", [
            ("Project classifier", "£0.50"),
            ("Review summariser", "£5"),
            ("Profile coach (LLM)", "£10"),
            ("Other (eval + dev)", "£15"),
        ], "~£30 / month"),
        ("Growth", "1,000 projects / month", [
            ("Project classifier", "£5"),
            ("Review summariser", "£40"),
            ("Profile coach (LLM)", "£80"),
            ("Decision assistant", "£60"),
        ], "~£200 / month"),
        ("Scale", "10,000 projects / month", [
            ("Project classifier", "£50"),
            ("Review summariser", "£200"),
            ("Profile coach (LLM)", "£300"),
            ("Decision assistant", "£250"),
        ], "~£800 / month"),
    ]

    tile_w = 4.0
    tile_h = 4.4
    gap = 0.25
    start_left = (SLIDE_W_IN - (tile_w * 3 + gap * 2)) / 2
    top = 2.3

    for i, (name, traffic, lines, total) in enumerate(bands):
        left = start_left + i * (tile_w + gap)
        panel(slide, left=left, top=top, width=tile_w, height=tile_h)
        # Header
        add_text(slide, text=name.upper(),
                 left=left, top=top + 0.25, width=tile_w, height=0.4,
                 size=11, bold=True, color=DIMMER, align=PP_ALIGN.CENTER)
        add_text(slide, text=traffic,
                 left=left, top=top + 0.6, width=tile_w, height=0.4,
                 size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        # Divider line
        line = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(left + 0.5), Inches(top + 1.15),
            Inches(left + tile_w - 0.5), Inches(top + 1.15),
        )
        line.line.color.rgb = PANEL_LIGHT
        line.line.width = Pt(0.75)

        # Line items
        y = top + 1.4
        for label, cost in lines:
            add_text(slide, text=label,
                     left=left + 0.4, top=y, width=tile_w - 1.4, height=0.35,
                     size=10, color=DIM)
            add_text(slide, text=cost,
                     left=left + 0.4, top=y, width=tile_w - 0.8, height=0.35,
                     size=10, color=TEXT, align=PP_ALIGN.RIGHT)
            y += 0.38

        # Total
        add_text(slide, text=total,
                 left=left, top=top + tile_h - 0.7, width=tile_w, height=0.45,
                 size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_text(slide,
             text="Cost is not the constraint. Engineering time and data quality are.",
             left=0.7, top=6.95, width=12, height=0.4,
             size=11, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 11)


def slide_12_synthetic_data(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Synthetic data")
    add_text(
        slide,
        text="Yes for some things. No for others. Don't confuse them.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )

    # 3-column matrix
    cols = [
        ("Pipeline plumbing",
         GREEN, "PERFECT",
         "Testing that data flows, jobs run, dashboards render, costs are sane.",
         ["Verify writes hit the right tables",
          "Run nightly crons end-to-end",
          "Profile latency + cost",
          "Catch integration bugs before launch"]),
        ("Prompt evaluation",
         GREEN, "GOOD with varied inputs",
         "Running an LLM against varied inputs to see if its outputs hold up.",
         ["Use LLM-generated test inputs",
          "Cover edge cases (typos, vague text)",
          "Manually inspect outputs",
          "Gate prompts on a checklist"]),
        ("Model training",
         ACCENT, "DANGEROUS",
         "Fitting a classifier or ranker that will go into production.",
         ["Model learns the SIM's biases, not reality",
          "Looks great in dev, collapses on real users",
          "Can't learn what isn't in the data",
          "ONE exception: fraud detection (slide 13)"]),
    ]

    col_w = 4.0
    col_h = 4.5
    gap = 0.2
    start_left = (SLIDE_W_IN - (col_w * 3 + gap * 2)) / 2
    top = 2.3

    for i, (name, color, verdict, body, points) in enumerate(cols):
        left = start_left + i * (col_w + gap)
        panel(slide, left=left, top=top, width=col_w, height=col_h, border=color)
        add_text(slide, text=name,
                 left=left + 0.25, top=top + 0.25, width=col_w - 0.5, height=0.4,
                 size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, text=verdict,
                 left=left + 0.25, top=top + 0.65, width=col_w - 0.5, height=0.35,
                 size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, text=body,
                 left=left + 0.25, top=top + 1.05, width=col_w - 0.5, height=0.7,
                 size=10, color=DIM, line_spacing=1.4)
        add_bullets(slide,
                    bullets=points,
                    left=left + 0.25, top=top + 1.95,
                    width=col_w - 0.5, height=2.4,
                    size=10, color=TEXT, bullet_char="·",
                    line_spacing=1.4)

    add_footer(slide, 12)


def slide_13_fraud_exception(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "The fraud exception", color=GREEN)
    add_text(
        slide,
        text="Fraud is the ONE place where synthetic data is\nbetter than real data.",
        left=0.7, top=1.0, width=12, height=1.5,
        size=26, bold=True,
    )

    add_text(
        slide,
        text=(
            "Fraud is a generative problem. Real fraudsters mechanically reproduce "
            "patterns. So can a sim. The sim's biases ARE the abuser's biases — and "
            "a model trained to detect mechanical patterns catches both."
        ),
        left=0.7, top=2.7, width=12, height=0.9,
        size=14, color=DIM, line_spacing=1.4,
    )

    # Recipe panel
    panel(slide, left=0.7, top=3.85, width=12, height=2.95)
    add_text(slide, text="THE RECIPE",
             left=0.95, top=4.0, width=11, height=0.35,
             size=11, bold=True, color=ACCENT)

    steps = [
        ("Step 1",
         "Write scripts/sim/seed-fraud-patterns.js",
         "Creates ~50 obvious sock-puppet clusters: shared IPs, fast signup-to-action latencies, identical user agents, templated comments, all liking the same builder."),
        ("Step 2",
         "Mark them as positive in abuse_signals",
         "Real users provide negative examples (legitimate accounts doing legitimate things). You now have a balanced labelled dataset on day 1."),
        ("Step 3",
         "Train an XGBoost classifier",
         "Inputs: signup IP, signup-to-first-action latency, device fingerprint cluster size, user-agent rarity, account age, posting velocity. Output: fraud probability 0–1."),
        ("Step 4",
         "Deploy and replace as real labels arrive",
         "Ship the classifier immediately. As admins moderate real flags, retrain on the union of sim + real data. The model improves continuously without you waiting."),
    ]

    y = 4.45
    for step, head, body in steps:
        add_text(slide, text=step,
                 left=0.95, top=y, width=1.0, height=0.4,
                 size=11, bold=True, color=ACCENT_DIM)
        add_text(slide, text=head,
                 left=2.05, top=y, width=10.7, height=0.32,
                 size=11, bold=True, color=TEXT)
        add_text(slide, text=body,
                 left=2.05, top=y + 0.3, width=10.7, height=0.4,
                 size=9, color=DIM, line_spacing=1.35)
        y += 0.6

    add_footer(slide, 13)


def slide_14_hybrid_pattern(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Sim → real graduation")
    add_text(
        slide,
        text="The hybrid pattern. Sim's role shifts as real data accumulates.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=24, bold=True,
    )

    phases = [
        ("Phase 1", "Today",
         "Sim only",
         "Build the entire AI pipeline using sim data. Verify writes, jobs, dashboards, costs all work. Do NOT deploy any trained models. Train rule-based or LLM-based features only."),
        ("Phase 2", "Weeks 4–8",
         "Sim for tests, real for production",
         "Production features run against real user data. Sim continues as the integration test environment — every CI run uses sim data to verify the pipeline still works. Start collecting outcomes in match_observations."),
        ("Phase 3", "Months 3–6",
         "Real for training, sim for backfill",
         "~500 real hire outcomes. Train the match engine on REAL outcomes only. Use sim data to backfill rare cases — e.g. simulate a project type you don't yet have real examples of, just to ensure the model doesn't crash on unseen inputs."),
        ("Phase 4", "Months 6+",
         "Real all the way",
         "Sim remains for tests + new feature dev only. All training data is real. Sim is no longer in the model training loop, only the QA loop."),
    ]

    y = 2.35
    for phase, when, role, body in phases:
        panel(slide, left=0.7, top=y, width=12, height=1.05)
        add_text(slide, text=phase,
                 left=0.95, top=y + 0.18, width=1.4, height=0.35,
                 size=14, bold=True, color=ACCENT)
        add_text(slide, text=when,
                 left=0.95, top=y + 0.55, width=1.4, height=0.3,
                 size=10, color=DIMMER)
        add_text(slide, text=role,
                 left=2.55, top=y + 0.13, width=10.0, height=0.35,
                 size=12, bold=True, color=ACCENT_DIM)
        add_text(slide, text=body,
                 left=2.55, top=y + 0.48, width=10.0, height=0.6,
                 size=10, color=TEXT, line_spacing=1.4)
        y += 1.15

    add_footer(slide, 14)


def slide_15_two_tricks(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Two underappreciated tricks")
    add_text(
        slide,
        text="LLMs make synthetic data dramatically more useful.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )

    tricks = [
        ("Trick 1",
         "Generate synthetic data WITH a stronger LLM",
         "Instead of your sim writing templated strings, ask Claude or GPT-4 to "
         "generate 200 varied recommendation comments / project descriptions / "
         "profile bios. Vary length, tone, completeness, formality, sentiment.\n\n"
         "Save the output as a fixture, sample from it in the sim. Suddenly your "
         "sim data is 10× more useful for prompt evaluation because it actually "
         "exposes your AI features to real-world diversity.",
         "Cost: ~£2 of API calls. Permanent improvement to your test environment."),

        ("Trick 2",
         "Use a stronger LLM to LABEL synthetic data",
         "If you ever want to train a small classifier without waiting for real "
         "data: generate synthetic inputs, then use a strong LLM to label them. "
         "The labels carry information your sim doesn't have.\n\n"
         "Train a small cheap classifier on those labels. Deploy the small model "
         "for production traffic. Standard distillation pattern — works "
         "particularly well for content quality and classification tasks.",
         "Saves 1000× on inference cost in production."),
    ]

    col_w = 6.0
    col_h = 4.55
    gap = 0.25
    start_left = (SLIDE_W_IN - (col_w * 2 + gap)) / 2
    top = 2.3

    for i, (label, head, body, footnote) in enumerate(tricks):
        left = start_left + i * (col_w + gap)
        panel(slide, left=left, top=top, width=col_w, height=col_h)
        add_text(slide, text=label.upper(),
                 left=left + 0.3, top=top + 0.25, width=col_w - 0.6, height=0.35,
                 size=11, bold=True, color=ACCENT)
        add_text(slide, text=head,
                 left=left + 0.3, top=top + 0.6, width=col_w - 0.6, height=0.6,
                 size=14, bold=True, color=TEXT, line_spacing=1.25)
        add_text(slide, text=body,
                 left=left + 0.3, top=top + 1.45, width=col_w - 0.6, height=2.5,
                 size=10, color=DIM, line_spacing=1.45)
        add_text(slide, text=footnote,
                 left=left + 0.3, top=top + col_h - 0.55, width=col_w - 0.6, height=0.4,
                 size=10, bold=True, color=ACCENT_DIM)

    add_footer(slide, 15)


def slide_16_week_one(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Week 1 plan", color=GREEN)
    add_text(
        slide,
        text="Day-by-day. Start Monday. Ship by Friday.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    days = [
        ("Day 1", "Schema",
         "Add the four new tables (project_classifications, match_observations, recommendation_signals, ai_inference_log). No code yet — just migration. Run in dev and prod."),
        ("Day 2", "Telemetry",
         "Wire up match_observations writes in the existing routes. Every render of the recommendations card → log it. Every hire → log the outcome. You're building the dataset."),
        ("Day 3", "Project classifier",
         "Single Express helper, single Anthropic prompt, structured-output JSON. Runs after every new project insert. Acceptance: post a fake project, see the classification land in the DB within 5 seconds."),
        ("Day 4", "Recommendation summariser",
         "Nightly cron route protected by an admin secret. For each builder with ≥3 recommendations, generate the 3-bullet summary. Cache. Display on the builder profile."),
        ("Day 5", "Profile coach (rule-based)",
         "Gap-detection rules on the tradesman dashboard. No AI yet — pure SQL. 'Add 3 more photos.' 'Verify with Companies House.' 'Your service area is small — consider expanding.'"),
    ]

    y = 2.4
    for day, name, body in days:
        panel(slide, left=0.7, top=y, width=12, height=0.85)
        add_text(slide, text=day,
                 left=0.95, top=y + 0.18, width=1.0, height=0.4,
                 size=14, bold=True, color=ACCENT)
        add_text(slide, text=name,
                 left=2.0, top=y + 0.13, width=3.0, height=0.4,
                 size=12, bold=True, color=ACCENT_DIM)
        add_text(slide, text=body,
                 left=5.1, top=y + 0.13, width=7.6, height=0.65,
                 size=10, color=TEXT, line_spacing=1.4)
        y += 0.95

    add_text(slide,
             text="End of week 1: 3 visible AI features in production, the dataset for everything else accumulating, and slide 8 of the pitch deck materially stronger.",
             left=0.7, top=7.0, width=12, height=0.4,
             size=11, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 16)


def slide_17_risks(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Risks & mitigations")
    add_text(
        slide,
        text="The four ways this goes wrong — and how we prevent each.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )

    risks = [
        ("Model trained on bad data",
         "Sim-trained matcher confidently makes wrong predictions on real users.",
         "Sim is for plumbing only. Train the matcher on real outcomes only. Hard rule."),
        ("Cost runaway",
         "An unbounded loop in a cron job racks up £1000 in LLM calls overnight.",
         "Cap monthly spend per feature in code. Log every call to ai_inference_log with cost. Alert at 50% budget."),
        ("LLM hallucinates on user content",
         "Summariser invents reviews that don't exist. Trust collapses.",
         "Always cite source recommendations in the summary. Show the original on hover. Cache aggressively + manual review queue for outliers."),
        ("Vendor lock-in to Anthropic",
         "Anthropic raises prices 5×. Or the model gets withdrawn.",
         "Abstract every LLM call behind a thin client interface. Swap providers in one config change. Keep prompts in version control."),
        ("Privacy / data handling",
         "User content sent to a third-party API. GDPR concerns.",
         "Anonymise PII before sending. Document the data flow. Add an opt-out flag for users who want strictly local processing."),
    ]

    y = 2.4
    for risk, problem, mitigation in risks:
        panel(slide, left=0.7, top=y, width=12, height=0.85)
        add_text(slide, text=risk,
                 left=0.95, top=y + 0.13, width=4.0, height=0.4,
                 size=12, bold=True, color=ACCENT_DIM)
        add_text(slide, text=problem,
                 left=0.95, top=y + 0.5, width=4.0, height=0.35,
                 size=9, color=DIM, line_spacing=1.3)
        add_text(slide, text=mitigation,
                 left=5.1, top=y + 0.18, width=7.6, height=0.65,
                 size=10, color=TEXT, line_spacing=1.4)
        y += 0.93

    add_footer(slide, 17)


def slide_18_success_metrics(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Success metrics")
    add_text(
        slide,
        text="How we know Project Lighthouse is working.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    milestones = [
        ("Week 1",
         "Pipeline live",
         "All four tables writing. 100+ projects classified. First builder summaries cached. Profile coach visible on tradesman dashboard."),
        ("Month 1",
         "Habits formed",
         "AI features used by every active homeowner. Builder coach delivers ≥1 actionable suggestion per profile. Average AI cost per project < £0.05."),
        ("Month 3",
         "First model trained",
         "≥500 closed hires logged in match_observations. Trust classifier trained on simulated + early real fraud labels. Caught 5+ true-positive fraud cases."),
        ("Month 6",
         "Match engine v2 live",
         "Match engine v2 in A/B test. Hire conversion rate from top-3 recommendations up by ≥10% vs deterministic baseline. Decision Assistant beta."),
        ("Month 12",
         "Moat in the pitch",
         "AI features cited in Series A pitch as the defensibility moat. Trust model has caught coordinated abuse competitors couldn't. >50% of homeowner hires flow through Decision Assistant."),
    ]

    y = 2.4
    for when, head, body in milestones:
        panel(slide, left=0.7, top=y, width=12, height=0.85)
        add_text(slide, text=when,
                 left=0.95, top=y + 0.2, width=1.4, height=0.35,
                 size=13, bold=True, color=ACCENT)
        add_text(slide, text=head,
                 left=2.5, top=y + 0.13, width=4.0, height=0.4,
                 size=12, bold=True, color=ACCENT_DIM)
        add_text(slide, text=body,
                 left=6.6, top=y + 0.18, width=6.1, height=0.6,
                 size=10, color=TEXT, line_spacing=1.4)
        y += 0.93

    add_footer(slide, 18)


# ─────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    print(f"building: {DECK_PATH.relative_to(ROOT)}")

    builders = [
        slide_01_title,
        slide_02_tldr,
        slide_03_why_now,
        slide_04_three_flywheels,
        slide_05_capability_map,
        slide_06_tier1_mvp,
        slide_07_tier2_data,
        slide_08_tier3_scale,
        slide_09_data_strategy,
        slide_10_stack_fit,
        slide_11_cost_model,
        slide_12_synthetic_data,
        slide_13_fraud_exception,
        slide_14_hybrid_pattern,
        slide_15_two_tricks,
        slide_16_week_one,
        slide_17_risks,
        slide_18_success_metrics,
    ]

    for i, build in enumerate(builders, start=1):
        print(f"  + slide {i:>2}  {build.__name__}")
        build(prs)

    DECK_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DECK_PATH))
    size_kb = os.path.getsize(DECK_PATH) // 1024
    print(f"\nsaved {DECK_PATH.relative_to(ROOT)}  ({size_kb} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()

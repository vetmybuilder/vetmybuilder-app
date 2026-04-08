#!/usr/bin/env python3
"""
build-pitch-deck.py

Populates the existing pitch-deck/vetmybuilder-pitch-deck-2026.pptx with
the 12 slides described in pitch-deck/slides.md, embedding the screenshots
from pitch-deck/screenshots/.

The script REPLACES any existing slide content. Run it freely after
re-shooting screenshots or updating the headlines.

Usage:
    python3 scripts/build-pitch-deck.py
"""

from __future__ import annotations

import copy
import os
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ── Paths ────────────────────────────────────────────────────────────

ROOT = Path(__file__).resolve().parents[1]
DECK_PATH = ROOT / "pitch-deck" / "vetmybuilder-pitch-deck-2026-v2.pptx"
SHOTS = ROOT / "pitch-deck" / "screenshots"
FLYWHEEL = ROOT / "pitch-deck" / "flywheel.png"
DEMO_MP4 = ROOT / "pitch-deck" / "demo.mp4"

# ── Theme ────────────────────────────────────────────────────────────

BG = RGBColor(0x0F, 0x17, 0x2A)        # slate-900
PANEL = RGBColor(0x1E, 0x29, 0x3B)     # slate-800
TEXT = RGBColor(0xF8, 0xFA, 0xFC)      # slate-50
DIM = RGBColor(0x94, 0xA3, 0xB8)       # slate-400
ACCENT = RGBColor(0xEF, 0x44, 0x44)    # red-500
ACCENT_DIM = RGBColor(0xFB, 0xA1, 0xA1)  # red-300

FONT = "Inter"  # falls back to Calibri / system default if not installed

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ── Helpers ──────────────────────────────────────────────────────────


def reset_deck(prs: Presentation) -> None:
    """Remove every existing slide so we always start clean."""
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001 — internal but stable
    rId_to_drop = []
    for sldId in list(sldIdLst):
        rId_to_drop.append(sldId.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"])
        sldIdLst.remove(sldId)
    for rId in rId_to_drop:
        prs.part.drop_rel(rId)


def find_blank_layout(prs: Presentation):
    """Return the first slide layout that's effectively blank."""
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name:
            return layout
    return prs.slide_layouts[6]  # PowerPoint convention: layout 6 = blank


def add_slide(prs: Presentation):
    layout = find_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    paint_background(slide)
    return slide


def paint_background(slide) -> None:
    """Cover the slide with a solid dark background rect."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False


def add_text(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_bullets(
    slide,
    *,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 20,
    color: RGBColor = TEXT,
    line_spacing: float = 1.35,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_image(
    slide,
    *,
    path: Path,
    left: float,
    top: float,
    width: float | None = None,
    height: float | None = None,
):
    if not path.exists():
        print(f"  ! missing image: {path.name}")
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def _add_arrowhead(connector):
    """
    Append a triangle arrowhead to the end of a connector. python-pptx
    doesn't expose this directly, so we modify the underlying line element.
    """
    line_elem = connector.line._get_or_add_ln()  # noqa: SLF001
    for old in line_elem.findall(qn("a:tailEnd")):
        line_elem.remove(old)
    tail = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


def add_eyebrow(slide, text: str, *, top: float = 0.55):
    add_text(
        slide,
        text=text.upper(),
        left=0.7,
        top=top,
        width=12,
        height=0.4,
        size=12,
        bold=True,
        color=ACCENT,
    )


def add_footer(slide, page_num: int, total: int = 18):
    add_text(
        slide,
        text="VetMyBuilder",
        left=0.7,
        top=SLIDE_H_IN - 0.55,
        width=4,
        height=0.3,
        size=10,
        color=DIM,
    )
    add_text(
        slide,
        text=f"{page_num} / {total}",
        left=SLIDE_W_IN - 1.3,
        top=SLIDE_H_IN - 0.55,
        width=1,
        height=0.3,
        size=10,
        color=DIM,
        align=PP_ALIGN.RIGHT,
    )


# ── Slide builders ───────────────────────────────────────────────────


def slide_1_title(prs: Presentation):
    slide = add_slide(prs)
    add_text(
        slide,
        text="VetMyBuilder",
        left=0.7,
        top=1.4,
        width=12,
        height=1.5,
        size=80,
        bold=True,
        color=TEXT,
    )
    add_text(
        slide,
        text="Trusted local builders, recommended by your neighbours.",
        left=0.7,
        top=2.7,
        width=12,
        height=0.8,
        size=24,
        color=DIM,
    )

    # Embedded demo video — single biggest engagement uplift on a cold deck
    if DEMO_MP4.exists():
        # add_movie expects a poster image; use the project-view hero as poster
        poster_path = SHOTS / "08-homeowner-project-view-with-recommendations.png"
        movie = slide.shapes.add_movie(
            str(DEMO_MP4),
            Inches(3.5),
            Inches(3.85),
            Inches(6.3),
            Inches(2.95),
            poster_frame_image=str(poster_path) if poster_path.exists() else None,
            mime_type="video/mp4",
        )
        add_text(
            slide,
            text="▶  Watch the 25-second product demo",
            left=3.5,
            top=6.85,
            width=6.3,
            height=0.4,
            size=12,
            color=ACCENT_DIM,
            align=PP_ALIGN.CENTER,
        )
    else:
        add_text(
            slide,
            text="[demo.mp4 missing — run scripts/record-pitch-demo.js]",
            left=3.5,
            top=4.5,
            width=6.3,
            height=0.5,
            size=14,
            color=DIM,
            align=PP_ALIGN.CENTER,
        )

    # Footer (no page number on title)
    add_text(
        slide,
        text="Pre-seed pitch  ·  [Month Year]  ·  [Founder name]",
        left=0.7,
        top=SLIDE_H_IN - 0.55,
        width=12,
        height=0.4,
        size=12,
        color=DIM,
    )


def slide_2_problem(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "The problem")
    add_text(
        slide,
        text="Finding a builder you can\ntrust is broken.",
        left=0.7,
        top=1.2,
        width=12,
        height=2.2,
        size=54,
        bold=True,
    )
    add_bullets(
        slide,
        bullets=[
            "1 in 3 UK homeowners regret who they hired.",
            "The average homeowner spends weeks searching before committing.",
            "Cowboy builders cost UK households billions every year.",
            "Word-of-mouth is what actually works — but it doesn’t scale.",
        ],
        left=0.7,
        top=4.3,
        width=12,
        height=2.6,
        size=22,
        color=DIM,
    )
    add_footer(slide, 2)


def slide_3_insight(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "The insight")
    add_text(
        slide,
        text="The most trusted recommendation\nis the one your neighbour just made.",
        left=0.7,
        top=2.3,
        width=12,
        height=2.5,
        size=44,
        bold=True,
    )
    add_text(
        slide,
        text=(
            "VetMyBuilder turns the trusted moments that already happen in WhatsApp\n"
            "groups, village Facebook pages and over the garden fence into a public,\n"
            "searchable, locality-anchored network."
        ),
        left=0.7,
        top=5.0,
        width=12,
        height=1.6,
        size=20,
        color=DIM,
    )
    add_footer(slide, 3)


def slide_4_solution(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "The solution")
    add_text(
        slide,
        text="Post a project. Your neighbours\nrecommend builders. You hire one.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.8,
        size=36,
        bold=True,
    )
    # Hero screenshot — the populated owner project view
    add_image(
        slide,
        path=SHOTS / "08-homeowner-project-view-with-recommendations.png",
        left=2.4,
        top=3.0,
        width=8.5,
    )
    add_footer(slide, 4)


def slide_5_homeowner_journey(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "How it works · homeowner side")
    add_text(
        slide,
        text="Three steps from idea to hired.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=36,
        bold=True,
    )

    captions = [
        ("1 · Post your project", "30 seconds. No email required to browse."),
        ("2 · Invite your network", "A magic link goes to neighbours and friends."),
        ("3 · Hire with confidence", "Ranked by who recommended, who liked, and accepted-hire history."),
    ]
    images = [
        SHOTS / "04-homeowner-signup-form.png",
        SHOTS / "10-builder-profile.png",
        SHOTS / "09-homeowner-project-view-hero.png",
    ]
    col_w = 4.0
    gap = 0.1
    start_left = (SLIDE_W_IN - (col_w * 3 + gap * 2)) / 2

    for i, (heading, body) in enumerate(captions):
        left = start_left + i * (col_w + gap)
        add_image(slide, path=images[i], left=left, top=2.4, width=col_w)
        add_text(
            slide,
            text=heading,
            left=left,
            top=5.45,
            width=col_w,
            height=0.4,
            size=16,
            bold=True,
            color=ACCENT_DIM,
        )
        add_text(
            slide,
            text=body,
            left=left,
            top=5.85,
            width=col_w,
            height=0.9,
            size=13,
            color=DIM,
        )

    add_footer(slide, 5)


def slide_6_tradesman_journey(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "How it works · tradesman side")
    add_text(
        slide,
        text="Leads earned, not bought.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=36,
        bold=True,
    )
    add_text(
        slide,
        text=(
            "Tradesmen create a free profile, see jobs in their area with the homeowner’s real ask, "
            "and accept hires that come with context an SMS lead never has — including who recommended them."
        ),
        left=0.7,
        top=1.95,
        width=12,
        height=0.9,
        size=15,
        color=DIM,
    )

    # Two screenshots side by side
    add_image(
        slide,
        path=SHOTS / "06-trade-signup-step1-company.png",
        left=0.7,
        top=3.0,
        width=6.0,
    )
    add_image(
        slide,
        path=SHOTS / "13-tradesman-own-profile.png",
        left=6.9,
        top=3.0,
        width=5.7,
    )
    add_footer(slide, 6)


def slide_7_flywheel(prs: Presentation):
    """
    Flywheel slide rendered with native PowerPoint shapes (rounded rects,
    oval, smart connectors) so the user can drag boxes around in PowerPoint
    and the arrows follow automatically.
    """
    slide = add_slide(prs)
    add_eyebrow(slide, "Why we win")
    add_text(
        slide,
        text="Every hire makes the next one easier.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=32,
        bold=True,
    )

    # Five node positions arranged clockwise around (6.67, 4.3).
    # Coordinates are slide-inch CENTRES — (cx, cy) — so the layout reads
    # naturally in source. Each card is 3.2 wide × 1.4 tall.
    nodes_def = [
        (
            "1 · Project posted",
            "A homeowner shares a job with their neighbours in 30 seconds.",
            6.67,
            1.95,
        ),
        (
            "2 · Neighbours recommend",
            "Trusted local builders are added by the people who used them.",
            10.95,
            3.6,
        ),
        (
            "3 · Tradesman accepts",
            "A real hire happens — the strongest possible trust signal.",
            9.4,
            5.95,
        ),
        (
            "4 · Builder score grows",
            "Every accepted hire compounds a builder's ranking on-platform.",
            3.93,
            5.95,
        ),
        (
            "5 · More homeowners trust",
            "Better signals → bigger network → even more projects.",
            2.38,
            3.6,
        ),
    ]

    card_w = 3.2
    card_h = 1.4
    nodes = []

    for heading, body, cx, cy in nodes_def:
        x = cx - card_w / 2
        y = cy - card_h / 2
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(card_w),
            Inches(card_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(2)
        rect.shadow.inherit = False

        tf = rect.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(72_000)
        tf.margin_right = Emu(72_000)
        tf.margin_top = Emu(36_000)
        tf.margin_bottom = Emu(36_000)

        # Heading paragraph
        p_head = tf.paragraphs[0]
        p_head.alignment = PP_ALIGN.CENTER
        run_head = p_head.add_run()
        run_head.text = heading
        run_head.font.size = Pt(13)
        run_head.font.bold = True
        run_head.font.color.rgb = TEXT
        run_head.font.name = FONT

        # Body paragraph
        p_body = tf.add_paragraph()
        p_body.alignment = PP_ALIGN.CENTER
        p_body.space_before = Pt(2)
        run_body = p_body.add_run()
        run_body.text = body
        run_body.font.size = Pt(10)
        run_body.font.color.rgb = DIM
        run_body.font.name = FONT

        nodes.append(rect)

    # Centre callout — solid oval with two-line label
    centre_size = 1.4
    centre_x = 6.67 - centre_size / 2
    centre_y = 4.3 - centre_size / 2
    centre = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(centre_x),
        Inches(centre_y),
        Inches(centre_size),
        Inches(centre_size),
    )
    centre.fill.solid()
    centre.fill.fore_color.rgb = PANEL
    centre.line.color.rgb = ACCENT
    centre.line.width = Pt(2)
    centre.shadow.inherit = False

    ctf = centre.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Emu(36_000)
    ctf.margin_right = Emu(36_000)
    ctf.margin_top = Emu(36_000)
    ctf.margin_bottom = Emu(36_000)

    cp_eyebrow = ctf.paragraphs[0]
    cp_eyebrow.alignment = PP_ALIGN.CENTER
    run_eyebrow = cp_eyebrow.add_run()
    run_eyebrow.text = "COMPOUNDING"
    run_eyebrow.font.size = Pt(9)
    run_eyebrow.font.bold = True
    run_eyebrow.font.color.rgb = ACCENT
    run_eyebrow.font.name = FONT

    cp_word = ctf.add_paragraph()
    cp_word.alignment = PP_ALIGN.CENTER
    run_word = cp_word.add_run()
    run_word.text = "trust"
    run_word.font.size = Pt(16)
    run_word.font.bold = True
    run_word.font.color.rgb = TEXT
    run_word.font.name = FONT

    # Smart connectors between the 5 nodes (clockwise loop).
    # Connection sites on a ROUNDED_RECTANGLE are: 0=top, 1=right, 2=bottom, 3=left.
    # Each tuple is (from_idx, from_site, to_idx, to_site).
    arrows_def = [
        (0, 1, 1, 0),  # 1 → 2  : right of node1  → top of node2
        (1, 2, 2, 0),  # 2 → 3  : bottom of node2 → top of node3
        (2, 3, 3, 1),  # 3 → 4  : left of node3   → right of node4
        (3, 0, 4, 2),  # 4 → 5  : top of node4    → bottom of node5
        (4, 1, 0, 3),  # 5 → 1  : right of node5  → left of node1
    ]

    for from_i, from_site, to_i, to_site in arrows_def:
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(0),
            Inches(0),
            Inches(0),
            Inches(0),
        )
        # begin_connect / end_connect link the connector to the shapes so
        # the user can drag the boxes around in PowerPoint and the arrows
        # automatically follow them.
        connector.begin_connect(nodes[from_i], from_site)
        connector.end_connect(nodes[to_i], to_site)
        connector.line.color.rgb = ACCENT
        connector.line.width = Pt(2.25)
        _add_arrowhead(connector)

    add_footer(slide, 7)


def slide_8_trust(prs: Presentation):
    """Defensibility / anti-abuse / trust system slide.

    Sits between the flywheel and the market slides — explains that the
    flywheel doesn't get gamed because the trust signals it produces are
    themselves protected by a layered anti-abuse system.
    """
    slide = add_slide(prs)
    add_eyebrow(slide, "Defensibility · trust system")
    add_text(
        slide,
        text="Trust by design, not by accident.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=32,
        bold=True,
    )
    add_text(
        slide,
        text=(
            "Marketplaces built on word-of-mouth live or die on whether they can stop "
            "fraud and self-dealing. We built an explicit, layered anti-abuse system "
            "from day one — not bolted on as a reaction."
        ),
        left=0.7,
        top=1.95,
        width=12,
        height=0.9,
        size=14,
        color=DIM,
    )

    # Left side: 4 pillars
    pillars = [
        (
            "Identity friction at signup",
            "Verified email, normalised duplicates (no Gmail dot-tricks), disposable-domain blocklist, age gates before first action.",
        ),
        (
            "Per-user trust score",
            "Every account scored 0–100 from age, verification, behaviour and graph signals. Likes and recommendations contribute proportionally — burner accounts add zero.",
        ),
        (
            "Behavioural fingerprints",
            "IP, device, fingerprint and email signals captured on every action. Cluster-detection cron flags accounts that look like the same person.",
        ),
        (
            "Live moderation queue",
            "Admin leaderboard surfaces suspicious patterns in real time. Cancelled hires, decline rates and audit-shadow rows are tracked and filterable.",
        ),
    ]
    pillar_left = 0.7
    pillar_top = 3.05
    pillar_w = 6.3
    row_h = 0.95
    for i, (head, body) in enumerate(pillars):
        y = pillar_top + i * row_h
        # Accent dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(pillar_left), Inches(y + 0.18), Inches(0.18), Inches(0.18)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        dot.shadow.inherit = False
        add_text(
            slide,
            text=head,
            left=pillar_left + 0.35,
            top=y,
            width=pillar_w,
            height=0.35,
            size=15,
            bold=True,
            color=ACCENT_DIM,
        )
        add_text(
            slide,
            text=body,
            left=pillar_left + 0.35,
            top=y + 0.35,
            width=pillar_w,
            height=0.6,
            size=11,
            color=DIM,
        )

    # Right side: admin leaderboard screenshot as proof
    add_image(
        slide,
        path=SHOTS / "14-admin-tradesmen-leaderboard.png",
        left=7.4,
        top=3.05,
        width=5.3,
    )
    add_text(
        slide,
        text="Live ranking signals on the internal admin leaderboard.",
        left=7.4,
        top=6.85,
        width=5.3,
        height=0.35,
        size=10,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, 8)


def slide_9_market(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "Market")
    add_text(
        slide,
        text="Big market. Bigger pain. Word-of-mouth wins.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=30,
        bold=True,
    )

    # Top row: 4 stat tiles. Numbers are real and citable — see slides.md.
    stats = [
        (
            "£11.3bn",
            "UK home improvement market\n(annual, 2024)",
            "IMARC Group",
        ),
        (
            "2m+",
            "Active tradespeople in the\nUK construction workforce",
            "ONS / Statista, Q3 2024",
        ),
        (
            "81%",
            "Homeowners who rely on personal\nrecommendations to find a tradesman",
            "Statista",
        ),
        (
            "£14.3bn",
            "Lost by UK households to\ncowboy builders in the last 5 years",
            "FMB, 2025",
        ),
    ]
    tile_w = 2.85
    tile_h = 2.05
    gap = 0.25
    start_left = (SLIDE_W_IN - (tile_w * 4 + gap * 3)) / 2
    top_y = 2.3

    for i, (big, label, source) in enumerate(stats):
        left = start_left + i * (tile_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top_y),
            Inches(tile_w),
            Inches(tile_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(1)
        rect.shadow.inherit = False

        add_text(
            slide,
            text=big,
            left=left,
            top=top_y + 0.18,
            width=tile_w,
            height=0.8,
            size=36,
            bold=True,
            color=ACCENT,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            text=label,
            left=left + 0.15,
            top=top_y + 1.0,
            width=tile_w - 0.3,
            height=0.7,
            size=11,
            color=DIM,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            text=source,
            left=left,
            top=top_y + 1.7,
            width=tile_w,
            height=0.3,
            size=9,
            color=DIM,
            align=PP_ALIGN.CENTER,
        )

    # Hero tile: the shadow market — strongest single number on the slide.
    # Filled red instead of bordered so it visually anchors the bottom half.
    hero_w = 9.5
    hero_h = 1.6
    hero_left = (SLIDE_W_IN - hero_w) / 2
    hero_top = 4.7

    hero_rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(hero_left),
        Inches(hero_top),
        Inches(hero_w),
        Inches(hero_h),
    )
    hero_rect.fill.solid()
    hero_rect.fill.fore_color.rgb = ACCENT
    hero_rect.line.fill.background()
    hero_rect.shadow.inherit = False

    add_text(
        slide,
        text="+ £10bn / year",
        left=hero_left,
        top=hero_top + 0.18,
        width=hero_w,
        height=0.6,
        size=30,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        text=(
            "UK economic activity lost every year because homeowners are too afraid to commission work — "
            "a shadow TAM that exists only because trust is broken."
        ),
        left=hero_left + 0.5,
        top=hero_top + 0.82,
        width=hero_w - 1.0,
        height=0.7,
        size=12,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )

    add_text(
        slide,
        text="Sources: IMARC Group · ONS Construction Statistics 2024 · Statista · Federation of Master Builders 2025",
        left=0.7,
        top=6.4,
        width=12,
        height=0.35,
        size=10,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        text=(
            "imarcgroup.com/uk-home-improvement-market   ·   "
            "ons.gov.uk/businessindustryandtrade/constructionindustry/articles/constructionstatistics/2024   ·   "
            "statista.com/statistics/428743   ·   "
            "fmb.org.uk/resource/cowboy-builders-could-be-costing-brits-a-staggering-14-3-billion.html   ·   "
            "fmb.org.uk/resource/consumer-fears-over-cowboy-builders-costs-uk-economy-10bn.html"
        ),
        left=0.7,
        top=6.7,
        width=12,
        height=0.4,
        size=7,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 10)


def slide_10_why_now(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "Why now")
    add_text(
        slide,
        text="Three things are changing at the same time.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=32,
        bold=True,
    )
    points = [
        ("1", "Trust in pay-per-lead marketplaces is collapsing.", "Tradesmen feel exploited; homeowners feel mis-sold. Trustpilot and Reddit tell the same story."),
        ("2", "Local networks are now digitally native.", "WhatsApp neighbour groups and village Facebook pages have already done the cultural work."),
        ("3", "AI lowered the cost of building it.", "A solo founder can now ship a marketplace MVP that needed a team of five three years ago."),
    ]
    col_w = 4.1
    gap = 0.2
    start_left = (SLIDE_W_IN - (col_w * 3 + gap * 2)) / 2

    for i, (num, head, body) in enumerate(points):
        left = start_left + i * (col_w + gap)
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(2.5), Inches(0.85), Inches(0.85)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        circle.shadow.inherit = False
        ctf = circle.text_frame
        ctf.margin_top = ctf.margin_bottom = ctf.margin_left = ctf.margin_right = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        crun.text = num
        crun.font.size = Pt(28)
        crun.font.bold = True
        crun.font.color.rgb = TEXT
        crun.font.name = FONT

        add_text(
            slide,
            text=head,
            left=left,
            top=3.5,
            width=col_w,
            height=1.3,
            size=18,
            bold=True,
        )
        add_text(
            slide,
            text=body,
            left=left,
            top=4.85,
            width=col_w,
            height=2.0,
            size=14,
            color=DIM,
        )

    add_footer(slide, 11)


def slide_11_business_model(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "Business model")
    add_text(
        slide,
        text="Free for homeowners. Revenue from\ntradesmen who win work.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.7,
        size=32,
        bold=True,
    )

    # Two columns: free side + paid side
    add_text(
        slide,
        text="Homeowner side",
        left=0.7,
        top=3.4,
        width=6,
        height=0.5,
        size=18,
        bold=True,
        color=ACCENT_DIM,
    )
    add_bullets(
        slide,
        bullets=[
            "Free, forever.",
            "Trust requires zero friction.",
            "No paywalled features.",
        ],
        left=0.7,
        top=3.9,
        width=6,
        height=2.5,
        size=16,
    )

    add_text(
        slide,
        text="Tradesman side",
        left=6.9,
        top=3.4,
        width=6,
        height=0.5,
        size=18,
        bold=True,
        color=ACCENT_DIM,
    )
    add_bullets(
        slide,
        bullets=[
            "Free tier — appear on shortlists, accept hires.",
            "Spotlight — featured placement on relevant projects.",
            "Contact unlocks — pay-per-quality-lead.",
            "Verified Pro (planned) — annual subscription with insurance & background checks.",
        ],
        left=6.9,
        top=3.9,
        width=6,
        height=3.0,
        size=15,
    )
    add_footer(slide, 12)


def slide_12_traction(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "Traction")
    add_text(
        slide,
        text="Where we are today.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=36,
        bold=True,
    )

    # 3 hero numbers
    numbers = [
        ("[N]", "homeowners signed up"),
        ("[N]", "recommendations made"),
        ("[N]", "hires accepted"),
    ]
    tile_w = 3.5
    tile_h = 2.6
    gap = 0.4
    start_left = (SLIDE_W_IN - (tile_w * 3 + gap * 2)) / 2
    for i, (big, label) in enumerate(numbers):
        left = start_left + i * (tile_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(2.6),
            Inches(tile_w),
            Inches(tile_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(1)
        rect.shadow.inherit = False
        add_text(
            slide,
            text=big,
            left=left,
            top=2.95,
            width=tile_w,
            height=1.3,
            size=64,
            bold=True,
            color=ACCENT,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            text=label,
            left=left,
            top=4.35,
            width=tile_w,
            height=0.6,
            size=15,
            color=DIM,
            align=PP_ALIGN.CENTER,
        )

    add_text(
        slide,
        text="“[Add a real customer quote here — it does emotional work the numbers can’t.]”",
        left=0.7,
        top=5.7,
        width=12,
        height=0.7,
        size=18,
        color=ACCENT_DIM,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        text="— [Customer name], [city]",
        left=0.7,
        top=6.4,
        width=12,
        height=0.4,
        size=12,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 16)


def slide_13_ask(prs: Presentation):
    slide = add_slide(prs)
    add_eyebrow(slide, "The ask")
    add_text(
        slide,
        text="We’re raising £350k to launch in Waltham Forest (E4–E17)\nand prove the recommendation flywheel.",
        left=0.7,
        top=1.0,
        width=12,
        height=2.0,
        size=26,
        bold=True,
    )

    add_text(
        slide,
        text="Use of funds",
        left=0.7,
        top=3.3,
        width=6,
        height=0.5,
        size=18,
        bold=True,
        color=ACCENT_DIM,
    )
    add_bullets(
        slide,
        bullets=[
            "45%  ·  £160k  ·  Engineering — first full-stack hire, 12 months",
            "25%  ·  £90k  ·  Marketing & community in Waltham Forest",
            "15%  ·  £50k  ·  Manual onboarding of the first 250 tradesmen",
            "15%  ·  £50k  ·  Founder runway, ops, legal, infra & buffer",
        ],
        left=0.7,
        top=3.85,
        width=6.0,
        height=2.7,
        size=13,
    )

    add_text(
        slide,
        text="What it gets us in 12 months",
        left=6.9,
        top=3.3,
        width=6,
        height=0.5,
        size=18,
        bold=True,
        color=ACCENT_DIM,
    )
    add_bullets(
        slide,
        bullets=[
            "2,000+ active homeowners across Waltham Forest",
            "300+ verified tradesmen onboarded",
            "250+ completed hires through the platform",
            "Series A readiness with K-factor > 1.2",
        ],
        left=6.9,
        top=3.85,
        width=6.0,
        height=2.7,
        size=13,
    )

    add_text(
        slide,
        text="Chris Morris   ·   [email]   ·   [phone]",
        left=0.7,
        top=6.65,
        width=12,
        height=0.4,
        size=14,
        color=ACCENT_DIM,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 18)


# ── New slides added in v3 (Competition / What's shipped / Team) ────


def slide_competition(prs: Presentation):
    """Pre-empts the 'what about Checkatrade' question."""
    slide = add_slide(prs)
    add_eyebrow(slide, "Competition")
    add_text(
        slide,
        text="Same market. Different incentives.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=32,
        bold=True,
    )

    rows = [
        (
            "",
            "Pay-per-lead\n(Checkatrade, Rated People)",
            "Review platforms\n(Trustpilot, Google)",
            "VetMyBuilder",
        ),
        (
            "Who pays",
            "Tradesmen pay per lead\nthey may not even win",
            "No one pays directly",
            "Tradesmen pay for visibility\nand contact unlocks",
        ),
        (
            "Trust signal",
            "Stranger reviews",
            "Stranger reviews",
            "Neighbours you actually know",
        ),
        (
            "Builder incentive",
            "Spend the most",
            "Game the reviews",
            "Deliver good work",
        ),
        (
            "Defensibility",
            "Brand and budget",
            "Scale",
            "Local recommendation graph\n+ trust system",
        ),
    ]

    grid_left = 0.7
    grid_top = 2.4
    col_widths = [2.4, 3.2, 3.2, 3.2]
    row_height = 0.85

    for r_idx, row in enumerate(rows):
        is_header = r_idx == 0
        y = grid_top + r_idx * row_height
        x = grid_left
        for c_idx, cell in enumerate(row):
            w = col_widths[c_idx]
            is_us = c_idx == 3

            if is_header or is_us:
                rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x),
                    Inches(y),
                    Inches(w),
                    Inches(row_height),
                )
                rect.fill.solid()
                if is_us and is_header:
                    rect.fill.fore_color.rgb = ACCENT
                elif is_us:
                    rect.fill.fore_color.rgb = PANEL
                else:
                    rect.fill.fore_color.rgb = PANEL
                if is_us and not is_header:
                    rect.line.color.rgb = ACCENT
                    rect.line.width = Pt(1.5)
                else:
                    rect.line.fill.background()
                rect.shadow.inherit = False

            if cell:
                color = TEXT if is_us or is_header else DIM
                bold = is_header or is_us or c_idx == 0
                add_text(
                    slide,
                    text=cell,
                    left=x + 0.15,
                    top=y + 0.1,
                    width=w - 0.3,
                    height=row_height - 0.2,
                    size=12 if not is_header else 13,
                    bold=bold,
                    color=color,
                    align=PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE,
                )
            x += w

    add_footer(slide, 9)


def slide_whats_shipped(prs: Presentation):
    """Solo-founder execution speed slide — your strongest non-traction asset."""
    slide = add_slide(prs)
    add_eyebrow(slide, "Momentum")
    add_text(
        slide,
        text="What we’ve shipped — solo, in weeks not months.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=30,
        bold=True,
    )
    add_text(
        slide,
        text=(
            "Most pre-seed pitches sell a roadmap. We sell what’s already running. "
            "Below is the build velocity of one founder over the last few months — "
            "all live in production or local dev today."
        ),
        left=0.7,
        top=1.95,
        width=12,
        height=0.9,
        size=14,
        color=DIM,
    )

    columns = [
        (
            "Marketplace core",
            [
                "Two-sided hire flow with multi-status tracking",
                "Recommendation system with magic-link sharing",
                "Builder profile + Companies House verification",
                "Tradesman wizard with strong-password checklist",
                "Spotlight + contact-unlock monetisation primitives",
            ],
        ),
        (
            "Trust + ranking",
            [
                "Recommendation scoring v2 — live, hire-weighted, recency-decayed",
                "Anti-abuse layered defence (identity, behaviour, graph)",
                "Admin moderation queue with audit-shadow filtering",
                "Trust score per user feeding the public ranking",
            ],
        ),
        (
            "Engineering quality",
            [
                "134 web unit tests, 51 server unit tests, all green",
                "End-to-end Playwright suite across 4 sharded shards",
                "Full local dev simulation with seeded multi-status data",
                "CI runs the entire suite on every push",
            ],
        ),
    ]
    col_w = 4.1
    gap = 0.2
    start_left = (SLIDE_W_IN - (col_w * 3 + gap * 2)) / 2

    for i, (heading, bullets) in enumerate(columns):
        left = start_left + i * (col_w + gap)
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(3.0),
            Inches(col_w),
            Inches(3.7),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL
        card.line.color.rgb = ACCENT
        card.line.width = Pt(1)
        card.shadow.inherit = False

        add_text(
            slide,
            text=heading,
            left=left + 0.25,
            top=3.15,
            width=col_w - 0.5,
            height=0.4,
            size=15,
            bold=True,
            color=ACCENT_DIM,
        )
        add_bullets(
            slide,
            bullets=bullets,
            left=left + 0.25,
            top=3.65,
            width=col_w - 0.5,
            height=3.0,
            size=11,
            color=TEXT,
            line_spacing=1.25,
        )

    add_footer(slide, 13)


def slide_ai_next(prs: Presentation):
    """
    Phase 2 — AI as the moat. Sits between "What we've shipped" and
    "Traction" so the narrative reads: here's what we've built → here's
    what we're building next → here's our traction.

    Investor framing of the Project Lighthouse internal plan: three
    flywheels, ship-now MVP, cost is a rounding error.
    """
    slide = add_slide(prs)
    add_eyebrow(slide, "What's next")
    add_text(
        slide,
        text="Phase 2 — AI is how we compound.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=32,
        bold=True,
    )
    add_text(
        slide,
        text=(
            "The marketplace is live. The next 12 months turn it into an "
            "AI-assisted service that gets smarter with every project, every "
            "recommendation, and every hire — and gives us a defensibility "
            "moat no competitor can replicate without running the same "
            "marketplace for the same length of time."
        ),
        left=0.7,
        top=1.95,
        width=12,
        height=0.95,
        size=14,
        color=DIM,
    )

    # Three flywheels
    flywheels = [
        (
            "Project understanding compounds",
            "Every project posted is a labelled training example. Within 1,000 projects we can train a classifier. Within 10,000, personalise it per project type and postcode.",
        ),
        (
            "Trust scores get sharper",
            "Every moderated abuse case trains the fraud model. After six months it catches novel attacks rule-based systems can't anticipate. The opposite of how most trust systems decay over time.",
        ),
        (
            "Match outcomes train the matcher",
            "Every closed hire is a label. The matcher learns which builder profiles win which kinds of projects in which postcodes. Unfakeable by competitors — they don't have the hire history.",
        ),
    ]

    card_w = 4.0
    card_h = 2.65
    gap = 0.2
    start_left = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2
    top = 3.1

    for i, (head, body) in enumerate(flywheels):
        left = start_left + i * (card_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(card_w),
            Inches(card_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(1)
        rect.shadow.inherit = False

        add_text(
            slide,
            text=f"0{i + 1}",
            left=left + 0.3,
            top=top + 0.25,
            width=card_w - 0.6,
            height=0.4,
            size=18,
            bold=True,
            color=ACCENT,
        )
        add_text(
            slide,
            text=head,
            left=left + 0.3,
            top=top + 0.7,
            width=card_w - 0.6,
            height=0.7,
            size=14,
            bold=True,
            color=TEXT,
        )
        add_text(
            slide,
            text=body,
            left=left + 0.3,
            top=top + 1.45,
            width=card_w - 0.6,
            height=1.1,
            size=11,
            color=DIM,
        )

    # Bottom callout — cost / shippability anchor.
    # (Phase 3 lives on its own dedicated slide now — see slide_phase3_mobile.)
    callout_left = 0.7
    callout_top = 6.05
    callout_w = 12
    callout_h = 0.85

    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(callout_left),
        Inches(callout_top),
        Inches(callout_w),
        Inches(callout_h),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = ACCENT
    callout.line.fill.background()
    callout.shadow.inherit = False

    add_text(
        slide,
        text=(
            "First four AI features ship in 4 weeks  ·  ~£30/month at MVP  ·  No new infrastructure  ·  Dataset starts compounding from day one"
        ),
        left=callout_left,
        top=callout_top + 0.22,
        width=callout_w,
        height=0.45,
        size=13,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, 14)


def slide_phase3_mobile(prs: Presentation):
    """
    Phase 3 — Project Toolbelt. Sits between slide 14 (Phase 2 — AI) and
    slide 16 (Traction). Mirrors slide 14's visual rhythm so the two
    pair as a "what's next, what's after that" sequence.
    """
    slide = add_slide(prs)
    add_eyebrow(slide, "Phase 3")
    add_text(
        slide,
        text="Project Toolbelt — the marketplace\nin every tradesman's pocket.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.7,
        size=30,
        bold=True,
    )
    add_text(
        slide,
        text=(
            "Tradesmen are on roofs and in vans, not at desks. The biggest "
            "weakness in any builder marketplace is response time on the "
            "supply side — and mobile push notifications fix it. Shipped "
            "after Phase 2, when we know which mobile features actually move "
            "the metric."
        ),
        left=0.7,
        top=2.55,
        width=12,
        height=0.95,
        size=14,
        color=DIM,
    )

    reasons = [
        (
            "Mobile-first humans",
            "A homeowner sits at a desk to research a renovation. A tradesman is on a roof. Web works for one side. Mobile works for the other.",
        ),
        (
            "Response time collapses",
            "Hire requests answered in minutes, not hours. Estimated 30–50% lift in acceptance rate within 90 days of launching the app.",
        ),
        (
            "Native where it matters",
            "Background push that works when the browser is closed. Camera-first photo capture for builder profiles. One-tap accept or decline.",
        ),
    ]

    card_w = 4.0
    card_h = 2.65
    gap = 0.2
    start_left = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2
    top = 3.7

    for i, (head, body) in enumerate(reasons):
        left = start_left + i * (card_w + gap)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(card_w),
            Inches(card_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(1)
        rect.shadow.inherit = False

        add_text(
            slide,
            text=f"0{i + 1}",
            left=left + 0.3,
            top=top + 0.25,
            width=card_w - 0.6,
            height=0.4,
            size=18,
            bold=True,
            color=ACCENT,
        )
        add_text(
            slide,
            text=head,
            left=left + 0.3,
            top=top + 0.7,
            width=card_w - 0.6,
            height=0.7,
            size=14,
            bold=True,
            color=TEXT,
        )
        add_text(
            slide,
            text=body,
            left=left + 0.3,
            top=top + 1.45,
            width=card_w - 0.6,
            height=1.1,
            size=11,
            color=DIM,
        )

    # Bottom callout — scope and timing
    callout_left = 0.7
    callout_top = 6.65
    callout_w = 12
    callout_h = 0.65

    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(callout_left),
        Inches(callout_top),
        Inches(callout_w),
        Inches(callout_h),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = ACCENT
    callout.line.fill.background()
    callout.shadow.inherit = False

    add_text(
        slide,
        text=(
            "PWA first (~1 month)  ·  React Native after (3–6 months)  ·  Tradesman-side scope only  ·  Triggered when Phase 2 ships"
        ),
        left=callout_left,
        top=callout_top + 0.15,
        width=callout_w,
        height=0.4,
        size=12,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, 15)


def slide_team(prs: Presentation):
    """Solo founder anchor slide. Pre-seed angels back people first."""
    slide = add_slide(prs)
    add_eyebrow(slide, "Team")
    add_text(
        slide,
        text="One founder. One mission.",
        left=0.7,
        top=1.0,
        width=12,
        height=1.0,
        size=36,
        bold=True,
    )

    # Left column: photo placeholder
    photo = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(2.6),
        Inches(3.2),
        Inches(3.8),
    )
    photo.fill.solid()
    photo.fill.fore_color.rgb = PANEL
    photo.line.color.rgb = ACCENT
    photo.line.width = Pt(1.5)
    photo.shadow.inherit = False
    add_text(
        slide,
        text="[ photo ]",
        left=0.9,
        top=4.2,
        width=3.2,
        height=0.6,
        size=14,
        color=DIM,
        align=PP_ALIGN.CENTER,
    )

    # Right column: bio
    add_text(
        slide,
        text="[Founder name]",
        left=4.5,
        top=2.65,
        width=8.2,
        height=0.6,
        size=28,
        bold=True,
    )
    add_text(
        slide,
        text="Founder & Engineer",
        left=4.5,
        top=3.2,
        width=8.2,
        height=0.4,
        size=16,
        color=ACCENT_DIM,
    )
    add_text(
        slide,
        text=(
            "[1–2 sentence bio. Background, prior shipped things, why this domain. "
            "Keep it concrete — investors skim. Example: ‘15 years building products at "
            "[company], shipped [thing] used by [N] people, got burned by a cowboy "
            "builder on my own extension last year and decided to fix it properly.’]"
        ),
        left=4.5,
        top=3.85,
        width=8.2,
        height=1.5,
        size=14,
        color=DIM,
    )

    # Three proof points
    proofs = [
        ("[N] years", "in [domain]"),
        ("Built", "the entire product solo"),
        ("Lives in", "Chingford — heart of\nWaltham Forest, your first market"),
    ]
    proof_top = 5.4
    proof_w = 2.65
    proof_h = 1.35
    proof_left_start = 4.5
    for i, (big, small) in enumerate(proofs):
        x = proof_left_start + i * (proof_w + 0.1)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(proof_top),
            Inches(proof_w),
            Inches(proof_h),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = PANEL
        rect.line.color.rgb = ACCENT
        rect.line.width = Pt(1)
        rect.shadow.inherit = False
        add_text(
            slide,
            text=big,
            left=x,
            top=proof_top + 0.12,
            width=proof_w,
            height=0.45,
            size=18,
            bold=True,
            color=ACCENT,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            text=small,
            left=x + 0.1,
            top=proof_top + 0.6,
            width=proof_w - 0.2,
            height=0.7,
            size=11,
            color=DIM,
            align=PP_ALIGN.CENTER,
        )

    add_footer(slide, 17)


# ── Main ─────────────────────────────────────────────────────────────


def main():
    if DECK_PATH.exists():
        print(f"opening: {DECK_PATH.relative_to(ROOT)}")
        prs = Presentation(str(DECK_PATH))
    else:
        print(f"creating fresh deck at: {DECK_PATH.relative_to(ROOT)}")
        prs = Presentation()
        DECK_PATH.parent.mkdir(parents=True, exist_ok=True)

    # Force 16:9 widescreen
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    print(f"resetting {len(prs.slides)} existing slide(s)")
    reset_deck(prs)

    builders = [
        slide_1_title,
        slide_2_problem,
        slide_3_insight,
        slide_4_solution,
        slide_5_homeowner_journey,
        slide_6_tradesman_journey,
        slide_7_flywheel,
        slide_8_trust,
        slide_competition,        # 9
        slide_9_market,           # 10
        slide_10_why_now,         # 11
        slide_11_business_model,  # 12
        slide_whats_shipped,      # 13
        slide_ai_next,            # 14  (Phase 2 — AI)
        slide_phase3_mobile,      # 15  (Phase 3 — Project Toolbelt)
        slide_12_traction,        # 16
        slide_team,               # 17
        slide_13_ask,             # 18
    ]

    for i, build in enumerate(builders, start=1):
        print(f"  + slide {i:>2}  {build.__name__}")
        build(prs)

    prs.save(str(DECK_PATH))
    size_kb = os.path.getsize(DECK_PATH) // 1024
    print(f"\nsaved {DECK_PATH.relative_to(ROOT)}  ({size_kb} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""
build-vault-deck.py

Builds the Project Vault internal financial deck — the cap table, equity
offer structure, dilution scenarios, burn-rate model, and term-sheet
playbook for VetMyBuilder's pre-seed raise.

INTERNAL ONLY. This is a reference document for the founder, not for
investors. Some pages contain sensitive cap-table arithmetic that should
never be in a deck handed to a third party.

Companion to:
    scripts/build-pitch-deck.py        (investor pitch)
    scripts/build-lighthouse-deck.py   (AI Phase 2 plan)

Usage:
    python3 scripts/build-vault-deck.py
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Paths ────────────────────────────────────────────────────────────

ROOT = Path(__file__).resolve().parents[1]
DECK_PATH = ROOT / "project-vault" / "project-vault-financials.pptx"

# ── Theme (matches the other decks) ─────────────────────────────────

BG = RGBColor(0x0F, 0x17, 0x2A)        # slate-900
PANEL = RGBColor(0x1E, 0x29, 0x3B)     # slate-800
PANEL_LIGHT = RGBColor(0x33, 0x40, 0x55)  # slate-700
TEXT = RGBColor(0xF8, 0xFA, 0xFC)      # slate-50
DIM = RGBColor(0x94, 0xA3, 0xB8)       # slate-400
DIMMER = RGBColor(0x64, 0x74, 0x8B)    # slate-500
ACCENT = RGBColor(0xEF, 0x44, 0x44)    # red-500
ACCENT_DIM = RGBColor(0xFB, 0xA1, 0xA1)  # red-300
GREEN = RGBColor(0x10, 0xB9, 0x81)     # emerald-500
AMBER = RGBColor(0xF5, 0x9E, 0x0B)     # amber-500

FONT = "Inter"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
TOTAL_SLIDES = 12

# ── Helpers ─────────────────────────────────────────────────────────


def find_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in (layout.name or "").lower():
            return layout
    return prs.slide_layouts[6]


def add_slide(prs: Presentation):
    slide = prs.slides.add_slide(find_blank_layout(prs))
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.2,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_bullets(
    slide,
    *,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 14,
    color: RGBColor = TEXT,
    line_spacing: float = 1.35,
    bullet_char: str = "•",
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_eyebrow(slide, text: str, *, top: float = 0.55, color: RGBColor = ACCENT):
    add_text(
        slide,
        text=text.upper(),
        left=0.7,
        top=top,
        width=12,
        height=0.4,
        size=12,
        bold=True,
        color=color,
    )


def add_footer(slide, page: int, total: int = TOTAL_SLIDES):
    add_text(
        slide,
        text="Project Vault  ·  Financials & Equity Plan  ·  INTERNAL",
        left=0.7,
        top=SLIDE_H_IN - 0.55,
        width=8,
        height=0.3,
        size=10,
        color=DIM,
    )
    add_text(
        slide,
        text=f"{page} / {total}",
        left=SLIDE_W_IN - 1.3,
        top=SLIDE_H_IN - 0.55,
        width=1,
        height=0.3,
        size=10,
        color=DIM,
        align=PP_ALIGN.RIGHT,
    )


def panel(slide, *, left: float, top: float, width: float, height: float,
          fill: RGBColor = PANEL, border: RGBColor = ACCENT,
          border_width: float = 1.0, rounded: bool = True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border is not None:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def hline(slide, *, x1, x2, y, color=PANEL_LIGHT, width=0.75):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y), Inches(x2), Inches(y),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


# ─────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────


def slide_01_title(prs):
    slide = add_slide(prs)
    add_text(
        slide,
        text="Project Vault",
        left=0.7, top=2.0, width=12, height=1.4,
        size=72, bold=True, color=TEXT,
    )
    add_text(
        slide,
        text="Financials & equity plan for the £350k pre-seed raise",
        left=0.7, top=3.35, width=12, height=0.7,
        size=22, color=ACCENT_DIM,
    )
    add_text(
        slide,
        text=(
            "An internal reference document — not to be sent to investors. Contains the cap "
            "table arithmetic, dilution scenarios, term-sheet playbook and burn model behind "
            "the headline numbers on slide 18 of the investor pitch."
        ),
        left=0.7, top=4.25, width=12, height=1.0,
        size=14, color=DIM,
    )

    pill = panel(
        slide, left=0.7, top=5.45, width=3.6, height=0.5,
        fill=PANEL, border=ACCENT, border_width=1.5,
    )
    add_text(
        slide,
        text="STATUS · DRAFT · INTERNAL",
        left=0.7, top=5.53, width=3.6, height=0.35,
        size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    add_text(
        slide,
        text="Companion to: vetmybuilder-pitch-deck-2026 · project-lighthouse-plan",
        left=0.7, top=SLIDE_H_IN - 0.65, width=12, height=0.4,
        size=11, color=DIM,
    )


def slide_02_tldr(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "TL;DR")
    add_text(
        slide,
        text="The deal in five numbers.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=36, bold=True,
    )

    stats = [
        ("£350k",     "Total raise",                   "Pre-seed"),
        ("17.5%",     "Equity offered",                "Standard UK pre-seed band"),
        ("£2.0m",     "Post-money valuation",          "£1.65m pre-money"),
        ("82.5%",     "Founder ownership after",       "Survives a Seed and Series A"),
        ("£250k+£100k", "SEIS + EIS split",            "Maximum UK tax incentive"),
    ]

    tile_w = 2.4
    tile_h = 3.4
    gap = 0.18
    start_left = (SLIDE_W_IN - (tile_w * 5 + gap * 4)) / 2
    top = 2.6

    for i, (big, label, sub) in enumerate(stats):
        left = start_left + i * (tile_w + gap)
        panel(slide, left=left, top=top, width=tile_w, height=tile_h)
        add_text(slide, text=big,
                 left=left, top=top + 0.55, width=tile_w, height=0.95,
                 size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        hline(slide, x1=left + 0.4, x2=left + tile_w - 0.4, y=top + 1.7)
        add_text(slide, text=label,
                 left=left + 0.2, top=top + 1.85, width=tile_w - 0.4, height=0.65,
                 size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(slide, text=sub,
                 left=left + 0.2, top=top + 2.55, width=tile_w - 0.4, height=0.6,
                 size=10, color=DIM, align=PP_ALIGN.CENTER)

    add_text(slide,
             text="Structure: Advance Subscription Agreement (ASA), £2.0m valuation cap, 12-month long-stop date.",
             left=0.7, top=6.4, width=12, height=0.4,
             size=12, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 2)


def slide_03_valuation_logic(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Valuation logic")
    add_text(
        slide,
        text="Why £2m post-money is defensible.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=30, bold=True,
    )

    add_text(
        slide,
        text=(
            "UK pre-seed valuation is driven mostly by founder profile + traction, not by financial metrics. "
            "Below is the standard band for solo technical founders in 2026, and where VetMyBuilder sits on it."
        ),
        left=0.7, top=1.95, width=12, height=0.7,
        size=13, color=DIM,
    )

    rows = [
        ("Solo non-tech founder · slides only",        "£500k – £1.0m",  False),
        ("Solo technical founder · prototype",          "£1.0m – £1.5m",  False),
        ("Solo technical founder · shipped product · no revenue", "£1.5m – £2.5m", True),
        ("Solo technical founder · shipped + early revenue", "£2.0m – £3.5m", False),
        ("Two-person team · shipped + early revenue",   "£2.5m – £4.0m",  False),
        ("Repeat founder · prior exit",                 "£3.0m – £6.0m+", False),
    ]

    y = 2.9
    row_h = 0.55
    for label, value, is_us in rows:
        if is_us:
            panel(slide, left=0.7, top=y, width=12, height=row_h, border=ACCENT, border_width=1.5)
            label_color = TEXT
            value_color = ACCENT
            label_bold = True
        else:
            label_color = DIM
            value_color = DIM
            label_bold = False

        add_text(slide, text=label,
                 left=0.95, top=y + 0.13, width=8.5, height=row_h - 0.2,
                 size=12, bold=label_bold, color=label_color)
        add_text(slide, text=value,
                 left=9.0, top=y + 0.13, width=3.7, height=row_h - 0.2,
                 size=12, bold=label_bold, color=value_color, align=PP_ALIGN.RIGHT)
        if is_us:
            add_text(slide, text="← VetMyBuilder",
                     left=9.0, top=y + 0.13, width=3.4, height=row_h - 0.2,
                     size=11, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
        y += row_h + 0.05

    add_text(slide,
             text="£1.65m pre-money / £2.0m post-money sits in the middle of the defensible band — neither desperate nor arrogant.",
             left=0.7, top=6.55, width=12, height=0.4,
             size=11, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 3)


def slide_04_cap_table(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Cap table")
    add_text(
        slide,
        text="Today vs after the pre-seed.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=32, bold=True,
    )

    # Two-column layout: BEFORE / AFTER
    columns = [
        ("Today",
         "Pre-fundraise",
         [
             ("Founder", "100.0%", "1,000,000 shares (notional)"),
         ],
         "100.0%"),
        ("After £350k pre-seed",
         "Post-money £2.0m · 17.5% offered",
         [
             ("Founder",   "82.5%", "1,000,000 shares"),
             ("Investors", "17.5%", "212,121 new shares issued"),
         ],
         "100.0%"),
    ]

    col_w = 6.0
    col_h = 4.7
    gap = 0.25
    start_left = (SLIDE_W_IN - (col_w * 2 + gap)) / 2
    top = 2.2

    for i, (title, sub, rows, total) in enumerate(columns):
        left = start_left + i * (col_w + gap)
        panel(slide, left=left, top=top, width=col_w, height=col_h)
        add_text(slide, text=title,
                 left=left + 0.4, top=top + 0.3, width=col_w - 0.8, height=0.5,
                 size=18, bold=True, color=ACCENT_DIM)
        add_text(slide, text=sub,
                 left=left + 0.4, top=top + 0.75, width=col_w - 0.8, height=0.35,
                 size=11, color=DIM)
        hline(slide, x1=left + 0.4, x2=left + col_w - 0.4, y=top + 1.2)

        y = top + 1.45
        for name, pct, note in rows:
            add_text(slide, text=name,
                     left=left + 0.4, top=y, width=2.5, height=0.4,
                     size=14, bold=True, color=TEXT)
            add_text(slide, text=pct,
                     left=left + 2.9, top=y, width=1.4, height=0.4,
                     size=14, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
            add_text(slide, text=note,
                     left=left + 0.4, top=y + 0.4, width=col_w - 0.8, height=0.35,
                     size=10, color=DIM)
            y += 0.95

        # Total
        hline(slide, x1=left + 0.4, x2=left + col_w - 0.4, y=top + col_h - 0.85)
        add_text(slide, text="TOTAL",
                 left=left + 0.4, top=top + col_h - 0.65, width=2.5, height=0.4,
                 size=12, bold=True, color=DIM)
        add_text(slide, text=total,
                 left=left + 2.9, top=top + col_h - 0.65, width=1.4, height=0.4,
                 size=12, bold=True, color=DIM, align=PP_ALIGN.RIGHT)

    add_footer(slide, 4)


def slide_05_dilution_scenarios(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Dilution scenarios")
    add_text(
        slide,
        text="The £350k cheque at five different equity prices.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )

    scenarios = [
        ("12%",  "£2.92m", "£2.57m", "Aggressive — only with warm leads",  False),
        ("15%",  "£2.33m", "£1.98m", "Confident — strong founder profile", False),
        ("17.5%", "£2.00m", "£1.65m", "Recommended — clean round number",   True),
        ("20%",  "£1.75m", "£1.40m", "Cautious — accept if no leads",      False),
        ("25%",  "£1.40m", "£1.05m", "Last resort — desperate signal",     False),
    ]

    # Table
    headers = ["Equity offered", "Post-money", "Pre-money", "Read"]
    col_widths = [2.4, 2.4, 2.4, 4.5]
    grid_left = (SLIDE_W_IN - sum(col_widths)) / 2
    top = 2.4
    row_h = 0.7

    # Header row
    x = grid_left
    for j, h in enumerate(headers):
        add_text(slide, text=h.upper(),
                 left=x + 0.15, top=top, width=col_widths[j] - 0.3, height=0.4,
                 size=10, bold=True, color=DIMMER,
                 align=PP_ALIGN.LEFT if j == 3 else PP_ALIGN.CENTER)
        x += col_widths[j]

    hline(slide, x1=grid_left + 0.15, x2=grid_left + sum(col_widths) - 0.15, y=top + 0.5)

    # Body rows
    y = top + 0.65
    for equity, post, pre, read, is_recommended in scenarios:
        if is_recommended:
            panel(slide, left=grid_left, top=y - 0.05, width=sum(col_widths), height=row_h, border=ACCENT, border_width=1.5)
            color = TEXT
            bold = True
        else:
            color = DIM
            bold = False

        x = grid_left
        cells = [equity, post, pre, read]
        for j, cell in enumerate(cells):
            add_text(slide, text=cell,
                     left=x + 0.15, top=y + 0.18, width=col_widths[j] - 0.3, height=0.4,
                     size=14 if j < 3 else 11,
                     bold=bold, color=color if j < 3 else (TEXT if is_recommended else DIM),
                     align=PP_ALIGN.LEFT if j == 3 else PP_ALIGN.CENTER)
            x += col_widths[j]
        y += row_h + 0.05

    add_text(slide,
             text="17.5% gives the cleanest post-money number (£2m) and leaves enough founder ownership for two more rounds.",
             left=0.7, top=6.5, width=12, height=0.4,
             size=11, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 5)


def slide_06_long_term_dilution(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Long-term dilution forecast")
    add_text(
        slide,
        text="Founder ownership through Seed and Series A.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )
    add_text(
        slide,
        text=(
            "Standard dilution per round: pre-seed 17.5%, seed 18%, Series A 20%, Series B 18%. "
            "Below: founder ownership trajectory under two scenarios."
        ),
        left=0.7, top=1.95, width=12, height=0.7,
        size=12, color=DIM,
    )

    # Two scenarios side by side
    scenarios = [
        (
            "Recommended — no pre-round option pool",
            "Negotiate the option pool to be created at the seed round, not pre-seed",
            [
                ("Today",                "100.0%"),
                ("After pre-seed (17.5%)", "82.5%"),
                ("After seed (~18%)",     "67.6%"),
                ("After Series A (~20%)", "54.1%"),
                ("After Series B (~18%)", "44.4%"),
            ],
            True,
        ),
        (
            "If forced into a 10% pre-round option pool",
            "Investor demands an option pool created BEFORE the round (dilutes you, not them)",
            [
                ("Today",                "100.0%"),
                ("After pre-seed + pool", "72.5%"),
                ("After seed (~18%)",     "59.5%"),
                ("After Series A (~20%)", "47.6%"),
                ("After Series B (~18%)", "39.0%"),
            ],
            False,
        ),
    ]

    col_w = 6.0
    col_h = 4.0
    gap = 0.25
    start_left = (SLIDE_W_IN - (col_w * 2 + gap)) / 2
    top = 2.85

    for i, (title, sub, rows, recommended) in enumerate(scenarios):
        left = start_left + i * (col_w + gap)
        border_color = ACCENT if recommended else PANEL_LIGHT
        border_w = 1.5 if recommended else 1.0
        panel(slide, left=left, top=top, width=col_w, height=col_h,
              border=border_color, border_width=border_w)

        add_text(slide, text=title,
                 left=left + 0.3, top=top + 0.2, width=col_w - 0.6, height=0.35,
                 size=13, bold=True, color=ACCENT_DIM if recommended else DIM)
        add_text(slide, text=sub,
                 left=left + 0.3, top=top + 0.55, width=col_w - 0.6, height=0.4,
                 size=10, color=DIMMER)

        hline(slide, x1=left + 0.3, x2=left + col_w - 0.3, y=top + 1.05)

        y = top + 1.2
        for label, pct in rows:
            add_text(slide, text=label,
                     left=left + 0.4, top=y, width=col_w - 1.6, height=0.35,
                     size=11, color=TEXT if recommended else DIM)
            color = ACCENT if (recommended and pct == "82.5%") else (TEXT if recommended else DIM)
            add_text(slide, text=pct,
                     left=left + col_w - 1.5, top=y, width=1.1, height=0.35,
                     size=11, bold=True, color=color, align=PP_ALIGN.RIGHT)
            y += 0.5

    add_footer(slide, 6)


def slide_07_seis_eis(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "SEIS / EIS structure")
    add_text(
        slide,
        text="The single most important UK-specific lever.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )
    add_text(
        slide,
        text=(
            "UK angels expect SEIS-eligible structures and most won't write a cheque without "
            "Advance Assurance from HMRC. Apply for it BEFORE you contact a single investor."
        ),
        left=0.7, top=1.95, width=12, height=0.7,
        size=13, color=DIM,
    )

    # Two cards: SEIS and EIS
    schemes = [
        ("SEIS",
         "Seed Enterprise Investment Scheme",
         "£250k",
         "First £250k of investment goes here",
         [
             "50% income tax relief for the investor",
             "Capital gains tax exemption on disposal",
             "Loss relief if the investment fails",
             "Investor risks ~£10k for every £20k cheque",
             "Maximum cap per company: £250k lifetime",
         ]),
        ("EIS",
         "Enterprise Investment Scheme",
         "£100k",
         "Used for the £100k above the SEIS cap",
         [
             "30% income tax relief for the investor",
             "Capital gains deferral",
             "No upper company cap until £12m raised",
             "Investor risks ~£14k for every £20k cheque",
             "Used for any amount beyond the SEIS £250k limit",
         ]),
    ]

    col_w = 6.0
    col_h = 4.1
    gap = 0.25
    start_left = (SLIDE_W_IN - (col_w * 2 + gap)) / 2
    top = 2.85

    for i, (name, full_name, amount, sub, points) in enumerate(schemes):
        left = start_left + i * (col_w + gap)
        panel(slide, left=left, top=top, width=col_w, height=col_h)

        add_text(slide, text=name,
                 left=left + 0.3, top=top + 0.2, width=col_w - 0.6, height=0.4,
                 size=20, bold=True, color=ACCENT)
        add_text(slide, text=full_name,
                 left=left + 0.3, top=top + 0.6, width=col_w - 0.6, height=0.35,
                 size=11, color=DIM)
        add_text(slide, text=amount,
                 left=left + col_w - 1.7, top=top + 0.2, width=1.4, height=0.4,
                 size=20, bold=True, color=TEXT, align=PP_ALIGN.RIGHT)
        add_text(slide, text=sub,
                 left=left + col_w - 3.2, top=top + 0.6, width=2.9, height=0.35,
                 size=10, color=DIM, align=PP_ALIGN.RIGHT)

        hline(slide, x1=left + 0.3, x2=left + col_w - 0.3, y=top + 1.1)

        add_bullets(slide,
                    bullets=points,
                    left=left + 0.35, top=top + 1.25,
                    width=col_w - 0.7, height=2.7,
                    size=11, color=TEXT, bullet_char="·",
                    line_spacing=1.4)

    add_footer(slide, 7)


def slide_08_use_of_funds(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Use of funds")
    add_text(
        slide,
        text="Where the £350k goes — line by line.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    items = [
        ("Engineering",        "£160k", "45%",
         "First full-stack hire for 12 months. London salary £90k–£100k base + employer NI + pension + tools = ~£100k–£120k all-in. Add £20k for laptop, dev tools, monitoring, hosting overflow, contractor backup."),
        ("Marketing & community", "£90k", "25%",
         "Two-sided GTM in Waltham Forest. ~£35k paid acquisition (Facebook + Google for homeowners at ~£8 CPL), £25k community + events + neighbour-group seeding, £15k freelance PR for 3 months around launch, £15k content + creative."),
        ("Supply-side acquisition", "£50k", "15%",
         "Concierge onboarding for the first 250 verified tradesmen. £200/tradesman: outreach + Companies House and insurance verification + activation incentives. Builds the supply side without overspending."),
        ("Founder runway, ops, legal, infra & buffer", "£50k", "15%",
         "Founder living costs £40k (£3.3k/month London). Legal setup, accountant, insurance, ICO, T&Cs, marketplace terms ~£10–15k. SaaS tools, hosting overflow ~£5–10k. Small contingency."),
    ]

    y = 2.4
    for name, amount, pct, body in items:
        panel(slide, left=0.7, top=y, width=12, height=1.0)
        add_text(slide, text=pct,
                 left=0.95, top=y + 0.18, width=1.2, height=0.6,
                 size=22, bold=True, color=ACCENT)
        add_text(slide, text=amount,
                 left=2.2, top=y + 0.22, width=1.5, height=0.4,
                 size=16, bold=True, color=ACCENT_DIM)
        add_text(slide, text=name,
                 left=2.2, top=y + 0.62, width=4.0, height=0.35,
                 size=11, color=TEXT)
        add_text(slide, text=body,
                 left=6.4, top=y + 0.15, width=6.4, height=0.85,
                 size=10, color=DIM, line_spacing=1.4)
        y += 1.1

    add_footer(slide, 8)


def slide_09_burn_runway(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Burn rate & runway")
    add_text(
        slide,
        text="~£29k/month average. 12 months of runway.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )
    add_text(
        slide,
        text=(
            "Burn ramps as the engineering hire onboards and the marketing budget kicks in. "
            "The £50k buffer absorbs ±10% slippage without changing the runway story."
        ),
        left=0.7, top=1.95, width=12, height=0.7,
        size=12, color=DIM,
    )

    phases = [
        ("Months 1–2",  "£18k / mo", "Setup phase",
         "Founder runway only. Legal incorporation, SEIS Advance Assurance, accountant, contracts. Engineer not hired yet."),
        ("Months 3–4",  "£28k / mo", "Onboarding phase",
         "Engineer joins. Marketing budget begins community outreach. Tradesman concierge starts in Waltham Forest."),
        ("Months 5–9",  "£32k / mo", "Full burn",
         "All four use-of-funds lines firing. Peak monthly burn during paid acquisition. The 'spending' months."),
        ("Months 10–12", "£28k / mo", "Wind-down + Series A prep",
         "Marketing tapers as organic K-factor takes over. Founder time shifts to fundraising for the next round."),
    ]

    card_w = 2.95
    card_h = 3.1
    gap = 0.18
    start_left = (SLIDE_W_IN - (card_w * 4 + gap * 3)) / 2
    top = 2.85

    for i, (when, burn, phase, body) in enumerate(phases):
        left = start_left + i * (card_w + gap)
        panel(slide, left=left, top=top, width=card_w, height=card_h)
        add_text(slide, text=when,
                 left=left + 0.25, top=top + 0.2, width=card_w - 0.5, height=0.35,
                 size=11, bold=True, color=DIMMER)
        add_text(slide, text=burn,
                 left=left + 0.25, top=top + 0.6, width=card_w - 0.5, height=0.55,
                 size=20, bold=True, color=ACCENT)
        add_text(slide, text=phase,
                 left=left + 0.25, top=top + 1.2, width=card_w - 0.5, height=0.35,
                 size=12, bold=True, color=ACCENT_DIM)
        hline(slide, x1=left + 0.25, x2=left + card_w - 0.25, y=top + 1.65)
        add_text(slide, text=body,
                 left=left + 0.25, top=top + 1.8, width=card_w - 0.5, height=1.2,
                 size=10, color=DIM, line_spacing=1.4)

    add_text(slide,
             text="Average burn: £29k/month  ·  Total: £350k  ·  Runway: 12 months  ·  Series A trigger: month 9",
             left=0.7, top=6.45, width=12, height=0.4,
             size=11, color=ACCENT_DIM, align=PP_ALIGN.CENTER)

    add_footer(slide, 9)


def slide_10_term_sheet(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Term sheet playbook")
    add_text(
        slide,
        text="What to accept. What to push back on.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )

    accept = [
        ("1× non-participating liquidation preference", "Standard"),
        ("Pro-rata rights for the lead",                 "Standard"),
        ("Information rights (quarterly updates)",       "Standard"),
        ("Weighted-average anti-dilution",               "Standard"),
        ("Standard ROFR / co-sale",                      "Standard"),
        ("Drag-along at 50%+ shareholder approval",      "Standard"),
    ]
    push_back = [
        ("Board seat at pre-seed",                       "Offer info rights only"),
        ("Pre-round option pool > 5%",                   "Create at seed round"),
        ("Full ratchet anti-dilution",                   "Use weighted-average"),
        ("Liquidation preference > 1× or participating", "Hold the line at 1× non-participating"),
        ("Founder vesting acceleration < 2× trigger",    "Single trigger on termination"),
        ("Drag-along < 50% threshold",                   "Never below 50%"),
    ]

    col_w = 6.0
    col_h = 4.4
    gap = 0.25
    start_left = (SLIDE_W_IN - (col_w * 2 + gap)) / 2
    top = 2.4

    cols = [
        ("ACCEPT",     "Standard terms — agree without arguing",   accept,    GREEN),
        ("PUSH BACK",  "Founder-killers — negotiate before signing", push_back, ACCENT),
    ]

    for i, (heading, sub, items, color) in enumerate(cols):
        left = start_left + i * (col_w + gap)
        panel(slide, left=left, top=top, width=col_w, height=col_h, border=color, border_width=1.5)
        add_text(slide, text=heading,
                 left=left + 0.3, top=top + 0.25, width=col_w - 0.6, height=0.4,
                 size=14, bold=True, color=color)
        add_text(slide, text=sub,
                 left=left + 0.3, top=top + 0.65, width=col_w - 0.6, height=0.35,
                 size=10, color=DIM)
        hline(slide, x1=left + 0.3, x2=left + col_w - 0.3, y=top + 1.05)

        y = top + 1.2
        for item, note in items:
            add_text(slide, text="·  " + item,
                     left=left + 0.35, top=y, width=col_w - 0.7, height=0.3,
                     size=10, color=TEXT)
            add_text(slide, text=note,
                     left=left + 0.55, top=y + 0.25, width=col_w - 0.85, height=0.25,
                     size=8, color=DIMMER)
            y += 0.5

    add_footer(slide, 10)


def slide_11_comparables(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Comparables")
    add_text(
        slide,
        text="What other UK pre-seed marketplaces have raised.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=26, bold=True,
    )
    add_text(
        slide,
        text=(
            "Public data is sparse for pre-seed because deals are private. The bands below are a "
            "rough composite from Crunchbase, Beauhurst, and SeedLegals public filings."
        ),
        left=0.7, top=1.95, width=12, height=0.7,
        size=12, color=DIM,
    )

    rows = [
        ("Marketplace pre-seed (UK, generic)",   "£150k – £400k",  "£1.0m – £2.5m post-money"),
        ("Marketplace with shipped product",     "£250k – £500k",  "£1.5m – £3.0m post-money"),
        ("Trades / home services marketplace",   "£300k – £750k",  "£2.0m – £4.0m post-money"),
        ("VetMyBuilder (target)",                "£350k",          "£2.0m post-money"),
        ("Marketplace with paying customers",    "£500k – £1.5m",  "£3.0m – £6.0m post-money"),
    ]

    col_widths = [5.8, 3.0, 3.5]
    grid_left = (SLIDE_W_IN - sum(col_widths)) / 2
    top = 2.95
    row_h = 0.6

    headers = ["Stage", "Round size", "Post-money"]
    x = grid_left
    for j, h in enumerate(headers):
        add_text(slide, text=h.upper(),
                 left=x + 0.15, top=top, width=col_widths[j] - 0.3, height=0.35,
                 size=10, bold=True, color=DIMMER)
        x += col_widths[j]
    hline(slide, x1=grid_left + 0.15, x2=grid_left + sum(col_widths) - 0.15, y=top + 0.45)

    y = top + 0.6
    for label, size, post in rows:
        is_us = "VetMyBuilder" in label
        if is_us:
            panel(slide, left=grid_left, top=y - 0.05, width=sum(col_widths), height=row_h, border=ACCENT, border_width=1.5)
        x = grid_left
        cells = [label, size, post]
        for j, cell in enumerate(cells):
            add_text(slide, text=cell,
                     left=x + 0.15, top=y + 0.13, width=col_widths[j] - 0.3, height=0.35,
                     size=12, bold=is_us,
                     color=ACCENT if is_us else DIM)
            x += col_widths[j]
        y += row_h + 0.05

    add_text(slide,
             text="Sources: Crunchbase pre-seed UK marketplace filter · Beauhurst SaaS reports · SeedLegals public deal data. Bands are indicative.",
             left=0.7, top=6.55, width=12, height=0.4,
             size=10, color=DIMMER, align=PP_ALIGN.CENTER)

    add_footer(slide, 11)


def slide_12_open_questions(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Open questions")
    add_text(
        slide,
        text="Decide before the first investor meeting.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=28, bold=True,
    )

    questions = [
        ("Apply for SEIS Advance Assurance?",
         "Yes — apply this week. Free, takes 4–8 weeks. UK angels expect it; without it, you cut your investor pool by ~80%."),
        ("ASA or priced round?",
         "ASA with a £2.0m cap and 12-month long-stop. Cheaper legal, faster to close, defers the exact dilution number until next round. Lawyers needed for either."),
        ("Lead investor or party round?",
         "Try for a lead first. A single £100k+ cheque from a recognised UK angel anchors the round and makes the rest fill in 2x faster. Party round is the fallback."),
        ("Option pool: pre-round or post-round?",
         "Push for post-round (created at seed). If forced into pre-round, cap at 5% maximum."),
        ("Founder vesting?",
         "Standard 4-year vest with 1-year cliff is normal. Refuse acceleration triggers below 'double-trigger' (change of control + termination)."),
        ("Cap-table software now or later?",
         "SeedLegals or Capdesk now — £40/month — saves a week of pain at the seed round when investors ask for a clean cap table."),
        ("Disclose competitive raises?",
         "Don't. UK pre-seed is a small world; mentioning a competing offer to one investor will reach others within 24 hours."),
    ]

    y = 2.4
    for q, a in questions:
        panel(slide, left=0.7, top=y, width=12, height=0.62)
        add_text(slide, text=q,
                 left=0.95, top=y + 0.1, width=4.5, height=0.45,
                 size=11, bold=True, color=ACCENT_DIM)
        add_text(slide, text=a,
                 left=5.55, top=y + 0.1, width=7.2, height=0.45,
                 size=10, color=TEXT, line_spacing=1.4)
        y += 0.7

    add_footer(slide, 12)


# ─────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    print(f"building: {DECK_PATH.relative_to(ROOT)}")

    builders = [
        slide_01_title,
        slide_02_tldr,
        slide_03_valuation_logic,
        slide_04_cap_table,
        slide_05_dilution_scenarios,
        slide_06_long_term_dilution,
        slide_07_seis_eis,
        slide_08_use_of_funds,
        slide_09_burn_runway,
        slide_10_term_sheet,
        slide_11_comparables,
        slide_12_open_questions,
    ]

    for i, build in enumerate(builders, start=1):
        print(f"  + slide {i:>2}  {build.__name__}")
        build(prs)

    DECK_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DECK_PATH))
    size_kb = os.path.getsize(DECK_PATH) // 1024
    print(f"\nsaved {DECK_PATH.relative_to(ROOT)}  ({size_kb} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()

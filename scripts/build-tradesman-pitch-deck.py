#!/usr/bin/env python3
"""
build-tradesman-pitch-deck.py

Generates a tradesman onboarding / sales pitch deck using the same visual
template as the investor pitch deck (dark slate-900 background, red accents,
Inter font, 16:9 widescreen).

Usage:
    python3 scripts/build-tradesman-pitch-deck.py
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# -- Paths -------------------------------------------------------------------

ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "pitch-deck"
OUT_PATH = OUT_DIR / "vetmybuilder-tradesman-pitch-deck.pptx"
QR_PATH = OUT_DIR / "qr-signup.png"

# -- Theme (identical to investor deck) --------------------------------------

BG = RGBColor(0x0F, 0x17, 0x2A)        # slate-900
PANEL = RGBColor(0x1E, 0x29, 0x3B)     # slate-800
TEXT = RGBColor(0xF8, 0xFA, 0xFC)      # slate-50
DIM = RGBColor(0x94, 0xA3, 0xB8)       # slate-400
ACCENT = RGBColor(0xEF, 0x44, 0x44)    # red-500
ACCENT_DIM = RGBColor(0xFB, 0xA1, 0xA1)  # red-300
GREEN = RGBColor(0x22, 0xC5, 0x5E)     # green-500

FONT = "Inter"

# -- Slide dimensions (16:9) -------------------------------------------------

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
TOTAL_SLIDES = 13

# -- Helpers (same API as investor deck) -------------------------------------


def find_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name:
            return layout
    return prs.slide_layouts[6]


def add_slide(prs: Presentation):
    layout = find_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    paint_background(slide)
    return slide


def paint_background(slide) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False


def add_text(
    slide,
    *,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_bullets(
    slide,
    *,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 20,
    color: RGBColor = TEXT,
    line_spacing: float = 1.35,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_eyebrow(slide, text: str, *, top: float = 0.55):
    add_text(
        slide,
        text=text.upper(),
        left=0.7,
        top=top,
        width=12,
        height=0.4,
        size=12,
        bold=True,
        color=ACCENT,
    )


def add_footer(slide, page_num: int, total: int = TOTAL_SLIDES):
    add_branded_logo(
        slide,
        left=0.7,
        top=SLIDE_H_IN - 0.55,
        width=4,
        height=0.3,
        size=10,
    )
    add_text(
        slide,
        text=f"{page_num} / {total}",
        left=SLIDE_W_IN - 1.3,
        top=SLIDE_H_IN - 0.55,
        width=1,
        height=0.3,
        size=10,
        color=DIM,
        align=PP_ALIGN.RIGHT,
    )


def add_panel(slide, *, left, top, width, height, border_color=ACCENT):
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = PANEL
    rect.line.color.rgb = border_color
    rect.line.width = Pt(1)
    rect.shadow.inherit = False
    return rect


def add_image(slide, *, path: Path, left: float, top: float, width: float | None = None, height: float | None = None):
    if not path.exists():
        print(f"  ! missing image: {path.name}")
        return None
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_branded_logo(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    """Render 'VetMyBuilder' with 'My' in accent red, rest in white."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align

    parts = [("Vet", TEXT), ("My", ACCENT), ("Builder", TEXT)]
    for text, color in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def add_number_circle(slide, num: str, left: float, top: float, size: float = 0.7):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TEXT
    r.font.name = FONT


# -- Slides ------------------------------------------------------------------


def slide_01_title(prs):
    slide = add_slide(prs)
    add_branded_logo(
        slide,
        left=0.7, top=1.4, width=12, height=1.5,
        size=80,
    )
    add_text(
        slide,
        text="Your reputation should be\nyour best marketing tool.",
        left=0.7, top=2.9, width=12, height=1.2,
        size=32, color=DIM,
    )
    add_text(
        slide,
        text="The trust-first platform for quality tradesmen.",
        left=0.7, top=4.4, width=12, height=0.6,
        size=20, color=ACCENT_DIM,
    )
    add_text(
        slide,
        text="Free to join  -  No lead fees  -  Built on trust, not ads",
        left=0.7, top=SLIDE_H_IN - 0.55, width=12, height=0.4,
        size=12, color=DIM,
    )


def slide_02_problem(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "The problem")
    add_text(
        slide,
        text="The platforms that were supposed\nto help you are working against you.",
        left=0.7, top=1.2, width=12, height=2.0,
        size=44, bold=True,
    )
    add_bullets(
        slide,
        bullets=[
            "You pay for leads that go to 5 other builders at the same time.",
            "You compete on price, not quality - it's a race to the bottom.",
            "Anonymous reviews can be faked or gamed by competitors.",
            "The same big companies always dominate the listings.",
            "The platforms profit whether you win the job or not.",
            "Your years of quality work are invisible online.",
        ],
        left=0.7, top=4.0, width=12, height=3.0,
        size=20, color=DIM,
    )
    add_footer(slide, 2)


def slide_03_shift(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "The shift")
    add_text(
        slide,
        text="Homeowners don't want more quotes.\nThey want someone they can trust.",
        left=0.7, top=1.8, width=12, height=2.0,
        size=44, bold=True,
    )
    add_text(
        slide,
        text=(
            "The way people choose builders is changing. They ask neighbours, friends,\n"
            "and family first - they always have. But until now, there's been no\n"
            "platform built around that."
        ),
        left=0.7, top=4.5, width=12, height=1.5,
        size=20, color=DIM,
    )
    add_footer(slide, 3)


def slide_04_introducing(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Introducing")
    add_text(
        slide,
        text="A platform built around how people\nactually choose builders.",
        left=0.7, top=1.0, width=12, height=1.6,
        size=36, bold=True,
    )
    add_text(
        slide,
        text=(
            "VetMyBuilder connects homeowners with trusted local tradesmen.\n"
            "Builders are ranked by trust signals: neighbour recommendations,\n"
            "verified credentials, quality of work, and responsiveness.\n\n"
            "No bidding wars. No pay-per-lead. No race to the bottom."
        ),
        left=0.7, top=2.8, width=12, height=2.0,
        size=20, color=DIM,
    )

    # Three key pillars as panel cards
    pillars = [
        ("Neighbour\nRecommendations", "Real people vouching\nfor real work"),
        ("Verified\nCredentials", "Companies House,\nGoogle, certifications"),
        ("Local\nMatching", "Homeowners in your area\nfind you first"),
    ]
    card_w = 3.6
    gap = 0.35
    start_left = (SLIDE_W_IN - (card_w * 3 + gap * 2)) / 2
    card_top = 5.1
    card_h = 1.6

    for i, (head, body) in enumerate(pillars):
        left = start_left + i * (card_w + gap)
        add_panel(slide, left=left, top=card_top, width=card_w, height=card_h)
        add_text(
            slide,
            text=head, left=left + 0.2, top=card_top + 0.15, width=card_w - 0.4, height=0.7,
            size=16, bold=True, color=ACCENT_DIM, align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            text=body, left=left + 0.2, top=card_top + 0.85, width=card_w - 0.4, height=0.6,
            size=12, color=DIM, align=PP_ALIGN.CENTER,
        )

    add_footer(slide, 4)


def slide_05_how_it_works(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "How it works")
    add_text(
        slide,
        text="Three steps to getting found\nby the right homeowners.",
        left=0.7, top=1.0, width=12, height=1.4,
        size=36, bold=True,
    )

    steps = [
        ("1", "Create your profile",
         "Company details, trade types, service\nareas, photos of your work,\ncertifications. Takes under 5 minutes."),
        ("2", "Get discovered",
         "Homeowners in your area see you.\nWhen neighbours recommend you,\nyour trust score grows."),
        ("3", "Win work on your terms",
         "Express interest in local projects.\nHomeowners shortlist based on trust,\nnot price. No competing against 10\nbuilders for the same lead."),
    ]
    col_w = 3.8
    gap = 0.3
    start_left = (SLIDE_W_IN - (col_w * 3 + gap * 2)) / 2

    for i, (num, head, body) in enumerate(steps):
        left = start_left + i * (col_w + gap)
        add_number_circle(slide, num, left + (col_w - 0.7) / 2, 2.8)
        add_text(
            slide,
            text=head, left=left, top=3.7, width=col_w, height=0.5,
            size=20, bold=True, color=ACCENT_DIM, align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            text=body, left=left, top=4.3, width=col_w, height=2.0,
            size=14, color=DIM, align=PP_ALIGN.CENTER,
        )

    add_footer(slide, 5)


def slide_06_different(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Why we're different")
    add_text(
        slide,
        text="This isn't another Checkatrade.",
        left=0.7, top=1.0, width=12, height=1.0,
        size=36, bold=True,
    )

    # Comparison table
    rows = [
        ("The Old Way", "The VetMyBuilder Way"),
        ("Pay per lead", "Free to join, pay only to boost"),
        ("Compete with 5+ builders per job", "Homeowners shortlist who they want"),
        ("Reviews anyone can fake", "Real neighbour recommendations"),
        ("Same big companies always on top", "Ranked by trust, not spend"),
        ("Generic, nationwide", "Hyperlocal, community-driven"),
        ("You're a lead", "You're a trusted professional"),
    ]

    table_left = 0.7
    table_top = 2.3
    col_w = 5.8
    row_h = 0.65
    gap = 0.35

    for i, (old, new) in enumerate(rows):
        y = table_top + i * row_h
        is_header = i == 0

        # Left column (old way)
        add_text(
            slide,
            text=old,
            left=table_left, top=y, width=col_w, height=row_h,
            size=14 if is_header else 18,
            bold=is_header,
            color=ACCENT if is_header else DIM,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Right column (VMB way)
        add_text(
            slide,
            text=new,
            left=table_left + col_w + gap, top=y, width=col_w, height=row_h,
            size=14 if is_header else 18,
            bold=is_header,
            color=GREEN if is_header else TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        if not is_header:
            # Subtle separator line
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(table_left), Inches(y), Inches(col_w * 2 + gap), Inches(0.01),
            )
            sep.fill.solid()
            sep.fill.fore_color.rgb = RGBColor(0x33, 0x3F, 0x55)
            sep.line.fill.background()
            sep.shadow.inherit = False

    add_footer(slide, 6)


def slide_07_recommendations(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Neighbour recommendations")
    add_text(
        slide,
        text="The most powerful marketing\nyou can't buy.",
        left=0.7, top=1.0, width=12, height=1.4,
        size=36, bold=True,
    )
    add_bullets(
        slide,
        bullets=[
            "When a homeowner recommends you, it's visible on your profile.",
            "This is the digital version of \"my neighbour used them and they were brilliant.\"",
            "Every job you do well becomes a trust signal for the next homeowner nearby.",
            "Your reputation compounds locally - the more you work in an area,\n    the stronger your presence.",
            "Neighbour recommendations carry more weight than anonymous online reviews.",
        ],
        left=0.7, top=2.8, width=12, height=3.5,
        size=20, color=DIM,
    )

    # Quote callout
    add_panel(slide, left=2.0, top=5.8, width=9.3, height=1.0)
    add_text(
        slide,
        text="\"Sarah in E4 recommended Acme Builders to 3 neighbours this month.\"",
        left=2.3, top=5.95, width=8.7, height=0.7,
        size=16, bold=False, color=TEXT, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_footer(slide, 7)


def slide_08_trust_signals(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Trust signals")
    add_text(
        slide,
        text="We verify what homeowners\nactually care about.",
        left=0.7, top=1.0, width=12, height=1.4,
        size=36, bold=True,
    )

    signals = [
        ("Neighbour recommendations", "Real people vouching for real work"),
        ("Companies House verification", "Confirmed registered business"),
        ("Google rating", "Pulled automatically, shown transparently"),
        ("Photos of completed work", "Real evidence, not stock images"),
        ("Trade types and specialisms", "Clear about what you do"),
        ("Service area", "Local to the homeowner"),
        ("Response time", "How quickly you engage"),
        ("Profile completeness", "Shows you take it seriously"),
    ]

    col_count = 2
    col_w = 5.5
    gap_x = 1.0
    row_h = 0.78
    start_left = (SLIDE_W_IN - (col_w * 2 + gap_x)) / 2
    start_top = 2.8

    for i, (head, body) in enumerate(signals):
        col = i % col_count
        row = i // col_count
        left = start_left + col * (col_w + gap_x)
        top = start_top + row * row_h

        # Accent dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top + 0.12), Inches(0.14), Inches(0.14)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        dot.shadow.inherit = False

        add_text(
            slide,
            text=head, left=left + 0.3, top=top, width=col_w - 0.3, height=0.35,
            size=16, bold=True, color=TEXT,
        )
        add_text(
            slide,
            text=body, left=left + 0.3, top=top + 0.35, width=col_w - 0.3, height=0.35,
            size=12, color=DIM,
        )

    add_text(
        slide,
        text="Every signal is earned, not bought. The more complete and verified your profile, the higher you rank.",
        left=0.7, top=6.3, width=12, height=0.5,
        size=14, color=ACCENT_DIM, align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 8)


def slide_09_pricing(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Pricing")
    add_text(
        slide,
        text="Free to join. Free to get found.\nPay only when you want more.",
        left=0.7, top=1.0, width=12, height=1.4,
        size=36, bold=True,
    )

    # Two-tier layout
    tiers = [
        (
            "Free - forever",
            GREEN,
            [
                "Create your full profile",
                "Be discovered by local homeowners",
                "Receive neighbour recommendations",
                "Companies House verification",
                "Google rating displayed",
                "Express interest in projects",
            ],
        ),
        (
            "Premium - coming soon",
            ACCENT,
            [
                "Spotlight placement in your area",
                "Priority matching for new projects",
                "Profile analytics and views",
                "Featured photos and portfolio",
                "Early access to new features",
                "Competitive pricing - a fraction of what lead-gen platforms charge",
            ],
        ),
    ]

    card_w = 5.5
    card_h = 3.8
    gap = 0.6
    start_left = (SLIDE_W_IN - (card_w * 2 + gap)) / 2
    card_top = 2.6

    for i, (title, title_color, items) in enumerate(tiers):
        left = start_left + i * (card_w + gap)
        add_panel(slide, left=left, top=card_top, width=card_w, height=card_h,
                  border_color=title_color)
        add_text(
            slide,
            text=title, left=left + 0.3, top=card_top + 0.2, width=card_w - 0.6, height=0.5,
            size=20, bold=True, color=title_color, align=PP_ALIGN.CENTER,
        )
        add_bullets(
            slide,
            bullets=items,
            left=left + 0.3, top=card_top + 0.85, width=card_w - 0.6, height=2.8,
            size=14, color=DIM, line_spacing=1.4,
        )

    add_text(
        slide,
        text=(
            "Your profile and visibility are always free. Premium features help you stand out more -\n"
            "they never gatekeep who can find you."
        ),
        left=0.7, top=6.6, width=12, height=0.6,
        size=13, color=DIM, align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 9)


def slide_10_early_adopter(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Early adopter advantage")
    add_text(
        slide,
        text="The builders who join first will always\nhave the strongest presence.",
        left=0.7, top=1.0, width=12, height=1.4,
        size=36, bold=True,
    )

    points = [
        "First movers get the most visibility in their area.",
        "You accumulate trust signals before anyone else.",
        "When homeowners start using VetMyBuilder, your profile is already established.",
        "Early adopters help shape the platform - your feedback shapes pricing and features.",
        "We onboard a limited number of quality builders per area to maintain standards.",
    ]
    add_bullets(
        slide,
        bullets=points,
        left=0.7, top=2.8, width=12, height=3.0,
        size=20, color=DIM,
    )

    # Urgency callout
    hero_w = 9.5
    hero_h = 1.2
    hero_left = (SLIDE_W_IN - hero_w) / 2
    hero_top = 5.6

    hero_rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(hero_left), Inches(hero_top), Inches(hero_w), Inches(hero_h),
    )
    hero_rect.fill.solid()
    hero_rect.fill.fore_color.rgb = ACCENT
    hero_rect.line.fill.background()
    hero_rect.shadow.inherit = False

    add_text(
        slide,
        text="This is the ground floor. The lift hasn't left yet.",
        left=hero_left, top=hero_top + 0.15, width=hero_w, height=0.4,
        size=22, bold=True, color=TEXT, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        text="We're rolling out area by area, starting in Waltham Forest.",
        left=hero_left, top=hero_top + 0.65, width=hero_w, height=0.4,
        size=14, color=TEXT, align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 10)


def slide_11_how_you_stand_out(prs):
    slide = add_slide(prs)
    add_eyebrow(slide, "Stand out")
    add_text(
        slide,
        text="How to build the strongest\nprofile on VetMyBuilder.",
        left=0.7, top=1.0, width=12, height=1.4,
        size=36, bold=True,
    )

    tips = [
        ("Complete your profile", "Add all trade types, service areas, and a company description. Completeness = credibility."),
        ("Upload photos of real work", "Kitchens, bathrooms, extensions - real projects, not stock images. This is your portfolio."),
        ("Get verified", "Companies House verification and your Google rating are pulled in automatically. The green tick matters."),
        ("Ask happy customers to recommend you", "When a homeowner you've worked for recommends you on VetMyBuilder, it's visible to their neighbours."),
        ("Respond quickly", "Response time is tracked. Quick, professional responses push you higher in rankings."),
        ("Stay local", "The more jobs you do in your area, the more your reputation compounds nearby."),
    ]

    col_w = 5.8
    gap_x = 0.5
    row_h = 0.8
    start_left = (SLIDE_W_IN - (col_w * 2 + gap_x)) / 2
    start_top = 2.8

    for i, (head, body) in enumerate(tips):
        col = i % 2
        row = i // 2
        left = start_left + col * (col_w + gap_x)
        top = start_top + row * (row_h + 0.45)

        add_text(
            slide,
            text=f"{i + 1}.",
            left=left, top=top, width=0.4, height=0.4,
            size=18, bold=True, color=ACCENT,
        )
        add_text(
            slide,
            text=head,
            left=left + 0.35, top=top, width=col_w - 0.35, height=0.35,
            size=16, bold=True, color=TEXT,
        )
        add_text(
            slide,
            text=body,
            left=left + 0.35, top=top + 0.35, width=col_w - 0.35, height=0.5,
            size=12, color=DIM,
        )

    add_footer(slide, 11)


def slide_12_vision(prs):
    slide = add_slide(prs)
    add_text(
        slide,
        text="We're building the most trusted way\nto find a builder in the UK.",
        left=0.7, top=1.8, width=12, height=2.0,
        size=44, bold=True, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        text=(
            "Starting local. Growing smart.\n"
            "Every feature rewards quality and transparency.\n\n"
            "The best builders don't need to shout.\n"
            "They just need to be seen."
        ),
        left=0.7, top=4.2, width=12, height=2.5,
        size=20, color=DIM, align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 12)


def slide_13_cta(prs):
    slide = add_slide(prs)
    add_text(
        slide,
        text="Ready to be found by\nthe right homeowners?",
        left=0.7, top=1.2, width=12, height=2.0,
        size=50, bold=True, align=PP_ALIGN.CENTER,
    )

    # CTA panel
    cta_w = 8.0
    cta_h = 3.9
    cta_left = (SLIDE_W_IN - cta_w) / 2
    cta_top = 3.5

    add_panel(slide, left=cta_left, top=cta_top, width=cta_w, height=cta_h, border_color=ACCENT)

    add_text(
        slide,
        text="Create your free profile in under 5 minutes.",
        left=cta_left, top=cta_top + 0.3, width=cta_w, height=0.5,
        size=22, bold=True, color=TEXT, align=PP_ALIGN.CENTER,
    )
    add_branded_logo(
        slide,
        left=cta_left, top=cta_top + 1.0, width=cta_w, height=0.5,
        size=28, align=PP_ALIGN.CENTER,
    )
    # QR code
    qr_size = 1.6
    qr_left = (SLIDE_W_IN - qr_size) / 2
    add_image(slide, path=QR_PATH, left=qr_left, top=cta_top + 1.6, width=qr_size, height=qr_size)

    add_text(
        slide,
        text="Scan to sign up  |  X  @vetmybuilder",
        left=cta_left, top=cta_top + 3.3, width=cta_w, height=0.4,
        size=14, color=DIM, align=PP_ALIGN.CENTER,
    )



# -- Main --------------------------------------------------------------------


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_shift(prs)
    slide_04_introducing(prs)
    slide_05_how_it_works(prs)
    slide_06_different(prs)
    slide_07_recommendations(prs)
    slide_08_trust_signals(prs)
    slide_09_pricing(prs)
    slide_10_early_adopter(prs)
    slide_11_how_you_stand_out(prs)
    slide_12_vision(prs)
    slide_13_cta(prs)

    prs.save(str(OUT_PATH))
    print(f"Saved {TOTAL_SLIDES}-slide tradesman pitch deck to {OUT_PATH}")


if __name__ == "__main__":
    main()
